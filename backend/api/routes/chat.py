"""
聊天 API 路由
"""
from fastapi import APIRouter, HTTPException, Body, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, Tuple, Dict, Any
from datetime import datetime, timedelta
from uuid import uuid4
import re
import json
import base64
import tempfile
import os
import hashlib
from pathlib import Path
from urllib.parse import urlparse

import asyncio

from loguru import logger
from config import settings, app_config
from database.init_db import conversations_db
from core.memory_manager import MemoryManager
from core.opencode_ws import OpenCodeClient, _conversation_current_session, is_conversation_running, mark_conversation_running, unmark_conversation_running
from core.memory_extractor import extract_and_save_background, extract_candidates
from core.growth import add_candidate, record_chat_growth_candidates
from core.skill_generator import generate_skill_body_from_chat
from core.skill_registry import (
    AUTO_GENERATED,
    BUILTIN,
    EXTERNAL,
    HERMES,
    OPENCLAW,
    OPENCODE,
    capture_opencode_skill_snapshot,
    get_skill_registry,
    migrate_new_opencode_skills_to_codebot,
    migrate_skill_dir_to_codebot_auto,
    opencode_skill_dirs,
    write_auto_skill_md,
)
from api.routes import scheduler as scheduler_router
from api.routes import mcp as mcp_router
from utils.installer import start_opencode_server

router = APIRouter()
opencode_ws: Optional[OpenCodeClient] = None
chat_memory_manager: Optional[MemoryManager] = None
# 由 main.py lifespan 注入（可为 None）
sandbox_manager = None


def _collect_opencode_retry_ports() -> List[int]:
    def _append_port(values: List[int], raw_port) -> None:
        try:
            port = int(str(raw_port).strip())
        except Exception:
            return
        if 1 <= port <= 65535 and port not in values:
            values.append(port)

    env_candidate_ports: List[int] = []
    _append_port(env_candidate_ports, os.environ.get("CODEBOT_OPENCODE_PREFERRED_PORT", ""))
    _append_port(env_candidate_ports, os.environ.get("CODEBOT_OPENCODE_FALLBACK_PORT", ""))
    if env_candidate_ports:
        return env_candidate_ports

    parsed = urlparse(app_config.opencode.server_url or "")
    configured_port = parsed.port or 11200
    candidate_ports: List[int] = []
    for raw_port in [configured_port, 11200, 4096, 50690]:
        _append_port(candidate_ports, raw_port)
    return candidate_ports


async def _ensure_opencode_client_connected(client: OpenCodeClient) -> bool:
    ok = await client.try_connect(attempts=3, delay=0.4, open_timeout=1.0)
    if ok:
        return True

    actual_port = 0
    for port in _collect_opencode_retry_ports():
        actual_port = await start_opencode_server(port)
        if actual_port:
            break

    if not actual_port:
        return False

    new_url = f"http://127.0.0.1:{actual_port}"
    client.base_url = new_url
    if opencode_ws is not None:
        opencode_ws.base_url = new_url
    return await client.try_connect(attempts=4, delay=0.4, open_timeout=1.0)

# 任务队列：key=conversation_id(str), value=asyncio.Queue[dict]
_task_queues: dict = {}
_queue_runners: dict = {}  # key=conversation_id -> asyncio.Task
_runtime_stream_state: Dict[str, Dict[str, Any]] = {}
_multi_agent_dispatch_state: Dict[str, Dict[str, Any]] = {}
_opencode_action_notification_keys: set[str] = set()


def _configured_reply_language() -> str:
    language = (getattr(app_config.general, "language", "") or "zh-CN").strip()
    return language if language in {"zh-CN", "en-US"} else "zh-CN"


def _growth_candidate_decision_enabled() -> bool:
    return bool(getattr(app_config.general, "growth_candidate_decision", False))


def _reply_language_instruction() -> str:
    if _configured_reply_language() == "en-US":
        return "Default to English for user-facing replies unless the user explicitly asks for another language."
    return "默认使用简体中文回复用户，除非用户明确要求使用其他语言。"


def _looks_like_skill_creation_intent(message: str) -> bool:
    text = re.sub(r"\s+", " ", (message or "").strip())
    if not text:
        return False
    lower = text.lower()
    if "skill-creator" in lower:
        return True
    has_skill_term = bool(re.search(r"(?<![a-z0-9_])skills?(?![a-z0-9_])|skill\.md|技能|能力", lower, flags=re.IGNORECASE))
    if not has_skill_term:
        return False
    use_only = re.search(r"(?:调用|使用|运行|执行|invoke|call|run|use)\s*(?:一个|这个|该)?\s*(?:skills?|技能)", lower, flags=re.IGNORECASE)
    create_near_skill = re.search(
        r"(?:创建|新建|生成|制作|写(?:一个|个)?|做成|沉淀|固化|保存(?:为|成)?|转(?:为|成)|提炼|提取|create|generate|make|write|save|materialize|scaffold)"
        r".{0,24}(?:skills?|skill\.md|技能|能力)",
        lower,
        flags=re.IGNORECASE,
    ) or re.search(
        r"(?:skills?|skill\.md|技能|能力).{0,24}"
        r"(?:创建|新建|生成|制作|写|做成|沉淀|固化|保存|转(?:为|成)|提炼|提取|create|generate|make|write|save|materialize|scaffold)",
        lower,
        flags=re.IGNORECASE,
    )
    return bool(create_near_skill and not use_only)


def _migrate_opencode_skill_refs_from_reply(assistant_response: str) -> List[Dict[str, Any]]:
    text = (assistant_response or "").replace("/", "\\")
    if not text:
        return []
    migrated: List[Dict[str, Any]] = []
    seen: set[str] = set()
    candidate_dirs: List[Path] = []

    for root in opencode_skill_dirs():
        root_text = str(root).replace("/", "\\")
        pattern = re.escape(root_text.rstrip("\\")) + r"\\+([^\\\r\n\"'`<>|]+)"
        for match in re.finditer(pattern, text, flags=re.IGNORECASE):
            slug = (match.group(1) or "").strip().strip(" .。；;，,")
            if slug:
                candidate_dirs.append(root / slug)

    for match in re.finditer(r"(?:^|[\\\s])\.agents\\skills\\([^\\\r\n\"'`<>|]+)", text, flags=re.IGNORECASE):
        slug = (match.group(1) or "").strip().strip(" .。；;，,")
        if slug:
            for root in opencode_skill_dirs():
                if root.name.lower() == "skills":
                    candidate_dirs.append(root / slug)

    for skill_dir in candidate_dirs:
        try:
            key = str(skill_dir.resolve()).lower()
        except Exception:
            key = str(skill_dir).lower()
        if key in seen:
            continue
        seen.add(key)
        item = migrate_skill_dir_to_codebot_auto(skill_dir, move=True)
        if not item:
            slug_base = re.sub(r"[^a-z0-9_]+", "_", skill_dir.name.strip().lower().replace("-", "_").replace(" ", "_"))
            slug_base = re.sub(r"_+", "_", slug_base).strip("_")
            auto_slug = f"auto_{slug_base}" if slug_base and not slug_base.startswith("auto_") else slug_base
            if auto_slug:
                registry = get_skill_registry()
                for existing_dir in settings.SKILLS_DIR.glob(f"{auto_slug}*"):
                    if existing_dir.is_dir() and (existing_dir / "SKILL.md").exists():
                        item = registry.find(f"auto:{existing_dir.name}")
                        if item:
                            break
        if item:
            item["migration_reason"] = "codebot_reply_skill_path"
            migrated.append(item)
    if migrated:
        logger.info(f"[skill] 已按助手回复中的 OpenCode skill 路径迁移 {len(migrated)} 个技能")
    return migrated


def _begin_codebot_skill_creation_watch(message: str) -> Optional[Dict[str, Any]]:
    if not _looks_like_skill_creation_intent(message):
        return None
    import time

    return {
        "snapshot": capture_opencode_skill_snapshot(),
        "since": time.time(),
        "message": message,
    }


def _finish_codebot_skill_creation_watch(context: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not context:
        return []
    migrated = migrate_new_opencode_skills_to_codebot(
        snapshot=context.get("snapshot") or {},
        since=context.get("since"),
        reason="codebot_chat_skill_creation",
    )
    if migrated:
        logger.info(f"[skill] Codebot 聊天已迁移 {len(migrated)} 个 OpenCode 生成技能到 Codebot skills")
    return migrated


class MessageRequest(BaseModel):
    """消息请求"""
    content: str
    conversation_id: Optional[int] = None


class MessageResponse(BaseModel):
    """消息响应"""
    success: bool
    data: dict
    message: str


class AttachedFile(BaseModel):
    name: str
    type: str   # mime type
    content: str  # base64 encoded for binary, plain text for text files
    is_text: bool = True


def _attachment_summary(file: AttachedFile) -> str:
    mime_type = str(file.type or "application/octet-stream")
    if mime_type.startswith("image/"):
        return f"【图片附件：{file.name}，类型：{mime_type}】"
    return f"【附件：{file.name}（二进制文件，类型：{mime_type}）】"


class SendMessageRequest(BaseModel):
    conversation_id: int
    message: str
    model: Optional[str] = None
    mode: Optional[str] = None  # agent mode: "plan", "build", or "agent"
    attached_files: Optional[List[AttachedFile]] = None
    user_already_saved: bool = False
    project_dir: Optional[str] = None  # 用户选择的项目文件夹路径
    target: Optional[str] = None  # codebot | hermes | obsidian
    knowledge_paths: Optional[List[str]] = None


class PermissionReplyRequest(BaseModel):
    request_id: str
    reply: str
    message: Optional[str] = None
    session_id: Optional[str] = None
    conversation_id: Optional[int] = None
    project_dir: Optional[str] = None


class QuestionReplyRequest(BaseModel):
    request_id: str
    answers: Optional[List[List[str]]] = None
    answer: Optional[str] = None
    reject: bool = False
    conversation_id: Optional[int] = None
    project_dir: Optional[str] = None

class UpdateTitleRequest(BaseModel):
    title: str

class TogglePinnedRequest(BaseModel):
    pinned: bool

class ToggleArchiveRequest(BaseModel):
    archived: bool

class ToggleGroupRequest(BaseModel):
    is_group: bool
    group_role: Optional[str] = None

class MultiAgentDispatchRequest(BaseModel):
    message: str
    model: Optional[str] = None
    mode: Optional[str] = None
    project_dir: Optional[str] = None

class ClearConversationRequest(BaseModel):
    confirm: bool = True

class SkillGenerateRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    version: str = "2.0.0"
    source: Optional[str] = None
    enabled: bool = True
    message_limit: int = 50

def generate_conversation_title(content: str, user_message: str = "") -> str:
    """本地回退标题生成：尽量提取任务主题，不直接照抄原始首句。"""
    max_length = 20

    _VERBOSE_PREFIXES = re.compile(
        r"^(?:"
        r"用户(?:想要?|需要|希望|要求|请求|说|问|提到|提出|想让我|让我|要我)"
        r"|(?:好的?|嗯|OK|Sure)[，,。.!！]?\s*(?:我来|让我|我帮你?|我将|下面我)"
        r"|(?:我来|让我|我帮你?|我将|下面我|我会|我可以)(?:帮你?|为你?|给你?)?"
        r"|(?:当然|没问题|可以|好嘞|好呀|好吧|好哒|行)[，,。.!！]?\s*"
        r"|(?:Hello|Hi|你好|您好)[，,。.!！]?\s*"
        r")\s*",
        re.IGNORECASE,
    )

    def _extract_title(raw: str) -> str:
        text = " ".join(raw.strip().split())
        text = re.sub(r"`{1,3}[\s\S]*?`{1,3}", "", text).strip()
        for sep in ["\n", "。", "！", "？", "；", ";", ".", "!", "?"]:
            if sep in text:
                text = text.split(sep, 1)[0].strip()
                break
        text = _VERBOSE_PREFIXES.sub("", text).strip()
        text = re.sub(r"^[\"'“”‘’《》【】\[\](){}]+|[\"'“”‘’《》【】\[\](){}]+$", "", text).strip()
        text = re.sub(r"^[，,。.!！？?；;：:\-—]+", "", text).strip()
        text = re.sub(r"(?:请问|帮我|麻烦|请你|我想|我要|如何|怎么|为什么)\s*", "", text).strip()
        if any(tok in text.lower() for tok in ["system_policy", "conversation_context", "internal_context"]):
            return ""
        return text

    # 优先从 AI 回复中提取标题（更能反映实际任务内容）
    if content:
        sanitized = _sanitize_assistant_output(content)
        title = _extract_title(sanitized)
        if title and len(title) >= 2:
            return title[:max_length]

    # 回退：从用户消息提取
    if user_message:
        title = _extract_title(user_message)
        if title and len(title) >= 2:
            return title[:max_length]

    return "新对话"


async def generate_conversation_title_via_ai(
    user_message: str,
    assistant_response: str,
    model: Optional[str] = None
) -> str:
    """调用 AI 根据用户问题和AI回复内容生成简洁的对话标题。
    
    AI 会综合理解对话处理的实际任务来命名，而非简单截取用户消息。
    如果 AI 调用失败，回退到本地 generate_conversation_title()。
    
    Args:
        model: 当前对话使用的模型（providerID/modelID 格式），传入以确保标题生成能正常工作。
    """
    max_length = 20
    try:
        client = opencode_ws
        if client is None or not getattr(client, "connected", False):
            logger.info(f"[标题生成] opencode_ws 未连接，尝试临时连接...")
            client = OpenCodeClient(app_config.opencode.server_url)
            ok = await _ensure_opencode_client_connected(client)
            if not ok:
                logger.warning("[标题生成] 无法连接 OpenCode，回退到本地生成")
                return generate_conversation_title(assistant_response, user_message=user_message)

        # 只截取少量上下文，尽快生成标题，避免额外浪费 token
        user_snippet = user_message[:220] if user_message else ""
        response_snippet = _sanitize_assistant_output(assistant_response or "", user_message=user_message)[:420]
        
        prompt = (
            f"请根据以下对话内容，生成一个简洁的中文对话标题（不超过{max_length}个字）。\n"
            "要求：\n"
            "1. 标题必须是对整个对话主题的总结，不要直接照抄用户提问第一句。\n"
            "2. 优先概括 AI 实际完成的任务、问题诊断结果或讨论主题。\n"
            "3. 不要输出引号、句号、冒号、序号、解释或多余文本。\n"
            "4. 只输出最终标题文本。\n\n"
            f"用户消息：{user_snippet}\n\n"
            f"AI回复：{response_snippet}"
        )
        logger.debug(f"[标题生成] 调用 AI 生成标题，model={model}")
        result = await client.execute_task(prompt, model=model, timeout=30)
        if result.success and result.content:
            title = " ".join(result.content.strip().split())
            title = title.split("\n", 1)[0].strip()
            title = title.strip('"\'“”‘’`')
            # 去除可能的多余标点和前缀
            title = re.sub(r"^(标题[：:]?\s*)", "", title).strip()
            title = re.sub(r"^[\-•*\d\.、\s]+", "", title).strip()
            title = re.sub(r"[。！？!?,，；;：:]+$", "", title).strip()
            if title and 2 <= len(title) <= max_length + 5:
                logger.info(f"[标题生成] AI 生成标题成功: {title}")
                return title[:max_length]
            else:
                logger.warning(f"[标题生成] AI 返回标题不合规（长度={len(title)}）: {title!r}")
        else:
            logger.warning(f"[标题生成] AI 调用失败: success={result.success}, error={result.error}")
    except Exception as e:
        logger.warning(f"[标题生成] AI 生成对话标题异常，回退本地生成: {e}")

    return generate_conversation_title(assistant_response, user_message=user_message)


async def _update_conversation_title_if_needed(
    conversation_id: int,
    user_message: str,
    assistant_response: str,
    model: Optional[str] = None,
) -> bool:
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            return False
        existing_title = (conversation.get("title") or "").strip()
        if existing_title and existing_title != "新对话":
            return False
        new_title = await generate_conversation_title_via_ai(user_message, assistant_response, model=model)
        if not new_title:
            return False
        await memory_manager.update_conversation_title(conversation_id, new_title)
        return True
    except Exception as e:
        logger.debug(f"后台更新对话标题失败: {e}")
        return False


def _assistant_excerpt_for_title(raw_content: str, user_message: str = "", min_chars: int = 180, max_chars: int = 420) -> str:
    cleaned = _sanitize_assistant_output(raw_content or "", user_message=user_message).strip()
    if len(cleaned) < min_chars:
        return ""
    return cleaned[:max_chars]


async def _run_chat_post_processing(
    user_message: str,
    assistant_response: str,
    conversation_id: Optional[int] = None,
    target: Optional[str] = None,
    execution_model: Optional[str] = None,
) -> None:
    skill_creation_intent = _looks_like_skill_creation_intent(user_message)
    migrated_skills: List[Dict[str, Any]] = []
    candidate_decision_enabled = _growth_candidate_decision_enabled()
    try:
        if candidate_decision_enabled:
            existing_contents: List[str] = []
            try:
                recent = await _get_chat_memory_manager().get_memories(limit=50)
                existing_contents = [str(item.get("content") or "") for item in recent if item.get("content")]
            except Exception:
                existing_contents = []
            candidates = await extract_candidates(
                user_message=user_message,
                assistant_response=assistant_response,
                existing_contents=existing_contents,
                memory_manager=_get_chat_memory_manager(),
                opencode_ws=opencode_ws,
            )
            for content, category in candidates:
                stage_memory_growth_candidate(
                    content=content,
                    category=category,
                    conversation_id=conversation_id,
                    evidence=user_message,
                )
        else:
            await extract_and_save_background(
                user_message=user_message,
                assistant_response=assistant_response,
                memory_manager=_get_chat_memory_manager(),
                opencode_ws=opencode_ws,
            )
    except Exception as exc:
        logger.debug(f"后台记忆提取失败: {exc}")

    try:
        await _stage_or_create_schedule_from_message(
            user_message,
            conversation_id=conversation_id,
            target=target,
            execution_model=execution_model,
        )
    except Exception as exc:
        logger.debug(f"后台定时任务沉淀失败: {exc}")

    if skill_creation_intent and not candidate_decision_enabled:
        try:
            migrated_skills.extend(_migrate_opencode_skill_refs_from_reply(assistant_response))
            if not migrated_skills:
                import time

                migrated_skills.extend(
                    migrate_new_opencode_skills_to_codebot(
                        snapshot=None,
                        since=time.time() - 600,
                        reason="codebot_post_processing_recent_skill",
                    )
                )
        except Exception as exc:
            logger.debug(f"后台 OpenCode skill 迁移兜底失败: {exc}")

    try:
        if not (skill_creation_intent and migrated_skills):
            await _stage_or_materialize_skill(
                user_message=user_message,
                assistant_response=assistant_response,
                conversation_id=conversation_id,
            )
    except Exception as exc:
        logger.debug(f"后台技能沉淀失败: {exc}")

    try:
        if not candidate_decision_enabled:
            candidates = record_chat_growth_candidates(
                user_message=user_message,
                assistant_response=assistant_response,
                conversation_id=conversation_id,
            )
            await _auto_apply_growth_candidates(
                candidates,
                user_message,
                assistant_response,
                conversation_id,
                target=target,
                execution_model=execution_model,
            )
    except Exception as exc:
        logger.debug(f"成长候选记录失败: {exc}")


async def _auto_apply_growth_candidates(
    candidates: List[Dict[str, Any]],
    user_message: str,
    assistant_response: str,
    conversation_id: Optional[int] = None,
    target: Optional[str] = None,
    execution_model: Optional[str] = None,
) -> None:
    """Close the learning loop for high-confidence or repeated growth candidates."""
    for candidate in candidates or []:
        try:
            if candidate.get("status") != "pending":
                continue
            confidence = float(candidate.get("confidence") or 0)
            hit_count = int(candidate.get("hit_count") or 1)
            kind = candidate.get("kind")
            candidate_id = str(candidate.get("id") or "")
            payload = candidate.get("payload") or {}
            should_apply = confidence >= 0.75 or hit_count >= 2
            if not should_apply or not candidate_id:
                continue

            if kind == "memory":
                await _get_chat_memory_manager().save_long_term_memory(
                    content=candidate.get("content") or user_message,
                    category=payload.get("category") or "note",
                    metadata={"source": "growth_auto_apply", "candidate_id": candidate_id},
                )
                from core.growth import ACCEPTED, mark_candidate
                mark_candidate(candidate_id, ACCEPTED)
            elif kind == "skill":
                source_user_message = payload.get("user_message") or user_message
                source_assistant_response = (
                    payload.get("assistant_response")
                    or candidate.get("content")
                    or assistant_response
                    or ""
                )
                item = await _materialize_skill_content(
                    name=candidate.get("title") or generate_conversation_title(user_message or "自动技能"),
                    description=(candidate.get("evidence") or user_message or "")[:180],
                    user_message=source_user_message,
                    assistant_response=source_assistant_response,
                )
                if item:
                    logger.info(f"[growth] 自动沉淀技能: {item.get('slug') or item.get('id')}")
                from core.growth import ACCEPTED, mark_candidate
                mark_candidate(candidate_id, ACCEPTED)
            elif kind == "task" and hit_count >= 2:
                executor = str(payload.get("executor") or "").strip() or _task_executor_from_target(target)
                model = str(payload.get("execution_model") or "").strip() or _task_execution_model_from_chat_model(execution_model)
                created = await _try_create_scheduled_task(candidate.get("content") or user_message, executor=executor, execution_model=model)
                if created:
                    from core.growth import ACCEPTED, mark_candidate
                    mark_candidate(candidate_id, ACCEPTED)
        except Exception as exc:
            logger.debug(f"成长候选自动落地失败: {exc}")


def stage_memory_growth_candidate(content: str, category: str = "note", conversation_id: Optional[int] = None, evidence: str = "") -> Optional[Dict[str, Any]]:
    content = str(content or "").strip()
    category = str(category or "note").strip() or "note"
    if not content:
        return None
    return add_candidate(
        kind="memory",
        title="待确认记忆",
        content=content,
        confidence=0.8,
        payload={"category": category, "conversation_id": conversation_id},
        evidence=str(evidence or content)[:1000],
    )


def stage_task_growth_candidate(
    title: str,
    content: str,
    *,
    conversation_id: Optional[int] = None,
    cron_expression: str = "",
    schedule_text: str = "",
    run_once: bool = False,
    executor: str = "opencode",
    execution_model: str = "",
    notify_channels: Optional[List[str]] = None,
    evidence: str = "",
) -> Optional[Dict[str, Any]]:
    title = str(title or "").strip()
    content = str(content or "").strip()
    if not title or not content:
        return None
    channels = notify_channels if isinstance(notify_channels, list) and notify_channels else ["app"]
    return add_candidate(
        kind="task",
        title=title,
        content=content,
        confidence=0.85 if cron_expression else 0.6,
        payload={
            "name": title,
            "task_prompt": content,
            "cron_expression": str(cron_expression or "").strip(),
            "schedule_text": str(schedule_text or "").strip(),
            "notify_channels": channels,
            "run_once": bool(run_once),
            "executor": _task_executor_from_target(executor),
            "execution_model": _task_execution_model_from_chat_model(execution_model),
            "conversation_id": conversation_id,
        },
        evidence=str(evidence or content)[:1000],
    )


async def notify_task_growth_candidate(candidate: Optional[Dict[str, Any]], conversation_id: Optional[int] = None) -> None:
    if not candidate:
        return
    if not bool(getattr(app_config.general, "task_candidate_notification_enabled", True)):
        return
    created_at = str(candidate.get("created_at") or "")
    updated_at = str(candidate.get("updated_at") or "")
    if created_at and updated_at and created_at != updated_at:
        return
    try:
        from api.routes import notifications as notifications_router

        service = getattr(notifications_router, "notification_service", None)
        if service is None:
            return
        title = "定时任务已进入成长候选"
        task_title = str(candidate.get("title") or "待确认定时任务")
        message = (
            f"{task_title}\n\n"
            "请打开右上角“成长候选”进行查看、编辑或接受。"
        )
        if conversation_id:
            message = f"{message}\n对话ID: {conversation_id}"
        asyncio.create_task(
            service.send_action_required_notification(
                title=title,
                message=message,
                task_id=f"growth:{candidate.get('id') or ''}",
                notif_type="warning",
                force_desktop=True,
            )
        )
    except Exception as exc:
        logger.debug(f"定时任务成长候选通知发送失败（跳过）: {exc}")


def _task_executor_from_target(target: Optional[str]) -> str:
    return "hermes" if str(target or "").strip().lower() == "hermes" else "opencode"


def _task_execution_model_from_chat_model(model: Optional[str]) -> str:
    return (
        str(model or "").strip()
        or app_config.general.chat_default_model
        or app_config.models.primary_model
        or ""
    ).strip()


def _task_executor_label(executor: Optional[str]) -> str:
    return "Hermes CLI" if _task_executor_from_target(executor) == "hermes" else "OpenCode CLI"


def _build_codebot_scheduler_boundary(executor: Optional[str] = None) -> str:
    return (
        "Codebot scheduler boundary:\n"
        "- The user's request is to create or update a Codebot scheduled task.\n"
        "- Do not create PowerShell background jobs, Windows schtasks, cron jobs, launchd jobs, systemd timers, or any other OS-level scheduler.\n"
        "- Do not execute the future task content immediately.\n"
        f"- Codebot will store the schedule in its own scheduler database; when it is due, Codebot will run it through {_task_executor_label(executor)} according to the task executor.\n"
        "- Reply only with the Codebot scheduling result or a concise acknowledgement."
    )


def _looks_like_codebot_schedule_creation_request(message: str) -> bool:
    return _looks_like_schedule_message(message or "")


async def _handle_codebot_schedule_creation_request(
    message: str,
    *,
    target: Optional[str] = None,
    conversation_id: Optional[int] = None,
    execution_model: Optional[str] = None,
) -> Optional[str]:
    if not _looks_like_codebot_schedule_creation_request(message):
        return None

    executor = _task_executor_from_target(target)
    model = _task_execution_model_from_chat_model(execution_model)
    result = await _stage_or_create_schedule_from_message(
        message,
        conversation_id=conversation_id,
        target=executor,
        execution_model=model,
    )
    if not result:
        result = (
            "我识别到这是定时任务创建请求，但 Codebot 没能解析出明确的执行时间，"
            "所以没有创建 PowerShell 后台作业、Windows 计划任务或其他系统级定时器。"
            "请补充具体时间，或到“定时任务”页面手动创建。"
        )
    return (
        f"{result}\n\n"
        f"执行器：{_task_executor_label(executor)}\n"
        f"执行模型：{model or '记忆整理备用模型'}\n"
        "调度方式：仅使用 Codebot 内置定时任务系统；未创建 PowerShell 后台作业、Windows schtasks 或其他系统级定时器。"
    )


def _default_task_notify_channels() -> List[str]:
    nc = app_config.notification
    channels: List[str] = []
    if nc.app_enabled:
        channels.append("app")
    if nc.desktop_enabled:
        channels.append("desktop")
    if nc.lark_enabled:
        channels.append("lark")
    if nc.email_enabled:
        channels.append("email")
    return channels or ["app"]


def _schedule_run_once_from_message(message: str) -> bool:
    text = message or ""
    return (
        any(key in text for key in ["今天", "明天", "后天", "一次", "只提醒一次", "仅提醒一次", "提醒一次"])
        or bool(re.search(r"\d+\s*(分钟|小时|天)\s*(后|之后|以后)", text))
        or bool(re.search(r"(半小时|一小时|一天)\s*(后|之后|以后)", text))
    )


async def _build_task_growth_candidate_payload(
    message: str,
    conversation_id: Optional[int] = None,
    target: Optional[str] = None,
    execution_model: Optional[str] = None,
) -> Optional[Dict[str, Any]]:
    raw_message = (message or "").strip()
    if not raw_message or not _looks_like_schedule_message(raw_message):
        return None

    notify_channels = _default_task_notify_channels()
    run_once = _schedule_run_once_from_message(raw_message)
    cron_expression = ""
    name = ""
    task_prompt = ""

    if _looks_like_birthday_reminder_message(raw_message):
        birthday = _extract_birthday_value(raw_message)
        if birthday:
            md = re.search(r"(\d{1,2})月(\d{1,2})日", birthday)
            if md:
                month = int(md.group(1))
                day = int(md.group(2))
                hour, minute = _extract_time_for_reminder(raw_message)
                subject = _extract_birthday_subject(raw_message) or "我"
                if subject == "我":
                    remind_text = "今天是你的生日，生日快乐！"
                    name = f"生日提醒：每年{month}月{day}日"
                else:
                    remind_text = f"今天是{subject}的生日，记得送上祝福。"
                    name = f"生日提醒：{subject}每年{month}月{day}日"
                task_prompt = f"__REMINDER__\n{remind_text}"
                cron_expression = f"{minute} {hour} {day} {month} *"

    if not cron_expression:
        try:
            cron_data = await scheduler_router.generate_cron_from_text(raw_message)
            cron_expression = str((cron_data or {}).get("cron") or "").strip()
        except Exception as exc:
            logger.debug(f"定时任务候选 cron 解析失败: {exc}")

    if not task_prompt:
        is_reminder = any(key in raw_message for key in ["提醒", "闹钟", "提示我", "叫我", "通知我", "别忘了", "记得"])
        if is_reminder:
            content = _extract_reminder_content(raw_message)
            name = f"{'一次性' if run_once else ''}提醒：{content}" if content else ("一次性定时提醒" if run_once else "定时提醒")
            task_payload = f"提醒：{content}" if content else raw_message
            task_prompt = f"__REMINDER__\n{task_payload}"
        else:
            task_content = _extract_task_content(raw_message)
            action_text = re.sub(r"[，。,.!！?？;；:：]", " ", task_content)
            action_text = " ".join(action_text.split()) or raw_message
            name = f"{'一次性' if run_once else ''}任务：{action_text}"
            task_prompt = task_content
            drive_match = re.search(r"保存到\s*([a-zA-Z])\s*盘", raw_message)
            if drive_match:
                drive = drive_match.group(1).upper()
                output_dir = f"{drive}:\\codebot_tasks"
                task_prompt = (
                    f"{task_content}\n"
                    f"请将产出保存为 Markdown 文件到 {output_dir} 目录（如不存在请创建），"
                    f"文件名包含日期时间（例如 20260301_0800.md），并在完成后输出保存路径。"
                )

    if len(name) > 30:
        name = f"{name[:30]}..."
    return {
        "title": name or "待确认定时任务",
        "content": task_prompt or raw_message,
        "cron_expression": cron_expression,
        "schedule_text": raw_message,
        "run_once": run_once,
        "executor": _task_executor_from_target(target),
        "execution_model": _task_execution_model_from_chat_model(execution_model),
        "notify_channels": notify_channels,
        "conversation_id": conversation_id,
    }


def stage_skill_growth_candidate(
    title: str,
    content: str,
    *,
    user_message: str,
    conversation_id: Optional[int] = None,
    evidence: str = "",
) -> Optional[Dict[str, Any]]:
    title = str(title or "").strip()
    content = _sanitize_assistant_output(content or "", user_message=user_message)
    if not title or not content:
        return None
    return add_candidate(
        kind="skill",
        title=title,
        content=content[:3000],
        confidence=0.8,
        payload={
            "user_message": user_message,
            "assistant_response": content[:6000],
            "conversation_id": conversation_id,
        },
        evidence=str(evidence or user_message or title)[:1000],
    )


async def _stage_or_create_schedule_from_message(
    message: str,
    conversation_id: Optional[int] = None,
    target: Optional[str] = None,
    execution_model: Optional[str] = None,
) -> Optional[str]:
    if _growth_candidate_decision_enabled():
        task_candidate = await _build_task_growth_candidate_payload(
            message,
            conversation_id=conversation_id,
            target=target,
            execution_model=execution_model,
        )
        if not task_candidate:
            return None
        staged = stage_task_growth_candidate(
            title=task_candidate.get("title") or "待确认定时任务",
            content=task_candidate.get("content") or message,
            conversation_id=conversation_id,
            cron_expression=task_candidate.get("cron_expression") or "",
            schedule_text=task_candidate.get("schedule_text") or message,
            run_once=bool(task_candidate.get("run_once")),
            executor=task_candidate.get("executor") or _task_executor_from_target(target),
            execution_model=task_candidate.get("execution_model") or _task_execution_model_from_chat_model(execution_model),
            notify_channels=task_candidate.get("notify_channels") if isinstance(task_candidate.get("notify_channels"), list) else None,
            evidence=message,
        )
        if not staged:
            return None
        await notify_task_growth_candidate(staged, conversation_id=conversation_id)
        return "已生成定时任务候选，可在“成长候选”中编辑后决定是否加入。"
    return await _try_create_scheduled_task(
        message,
        executor=_task_executor_from_target(target),
        execution_model=_task_execution_model_from_chat_model(execution_model),
    )


async def _stage_or_materialize_skill(user_message: str, assistant_response: str, conversation_id: Optional[int] = None) -> bool:
    if _growth_candidate_decision_enabled():
        if not _should_materialize_skill(user_message, assistant_response, conversation_id=conversation_id):
            return False
        staged = stage_skill_growth_candidate(
            title=generate_conversation_title(user_message or "自动技能"),
            content=assistant_response,
            user_message=user_message,
            conversation_id=conversation_id,
            evidence=user_message,
        )
        return bool(staged)
    return await _materialize_reusable_skill(
        user_message=user_message,
        assistant_response=assistant_response,
        conversation_id=conversation_id,
    )


async def _materialize_skill_content(
    *,
    name: str,
    description: str,
    user_message: str,
    assistant_response: str,
    slug_hint: str = "",
) -> Optional[Dict[str, Any]]:
    cleaned_response = _sanitize_assistant_output(assistant_response or "", user_message=user_message)
    if _skill_content_is_noise(cleaned_response):
        return None

    new_name = str(name or "").strip() or f"自动技能-{generate_conversation_title(user_message or '自动技能')}"
    desc = str(description or "").strip() or cleaned_response.strip().replace("\n", " ")
    if _skill_content_is_noise(desc):
        return None
    if len(desc) > 180:
        desc = f"{desc[:180]}..."

    try:
        all_skills = get_skill_registry().list_skills()
    except Exception:
        all_skills = []

    similar_skill = None
    for skill in all_skills:
        existing_name = skill.get("name", "")
        existing_desc = skill.get("description", "")
        if _skill_name_similar(existing_name, new_name) or _skill_name_similar(existing_desc, desc):
            similar_skill = skill
            break

    if similar_skill:
        skill_id = similar_skill.get("id", "")
        if similar_skill.get("source") == AUTO_GENERATED or skill_id.startswith("auto:"):
            slug = skill_id.split(":", 1)[1].strip()
            skill_md = settings.SKILLS_DIR / slug / "SKILL.md"
            if skill_md.exists():
                skill_body = await generate_skill_body_from_chat(
                    user_message=user_message,
                    assistant_response=cleaned_response,
                    title=new_name,
                    description=desc,
                    opencode_client=opencode_ws if opencode_ws and getattr(opencode_ws, "connected", False) else None,
                )
                _write_auto_skill_md(skill_md, new_name, desc, user_message, skill_body)
                logger.info(f"升级已有自动技能: {slug}")
                return get_skill_registry().find(f"auto:{slug}") or similar_skill
        logger.info(
            f"[skill] 跳过技能生成：已有相似技能 '{similar_skill.get('name')}' "
            f"(id={skill_id})"
        )
        return similar_skill

    digest = hashlib.sha1(f"{user_message}\n{cleaned_response[:300]}".encode("utf-8")).hexdigest()[:8]
    slug_base = re.sub(r"[^a-z0-9_]", "_", (slug_hint or new_name)[:20].lower().replace(" ", "_"))
    slug_base = re.sub(r"_+", "_", slug_base).strip("_") or "auto"
    skill_body = await generate_skill_body_from_chat(
        user_message=user_message,
        assistant_response=cleaned_response,
        title=new_name,
        description=desc,
        opencode_client=opencode_ws if opencode_ws and getattr(opencode_ws, "connected", False) else None,
    )
    item = get_skill_registry().create_auto_skill(
        name=new_name,
        description=desc,
        body=skill_body,
        user_message=user_message,
        slug_hint=f"{slug_base}_{digest}",
    )
    logger.info(f"生成新自动技能（本地）: {item.get('slug')}")
    return item


def _sanitize_assistant_output(content: str, user_message: str = "") -> str:
    if not content:
        return ""
    text = str(content).replace("\r\n", "\n")

    # ── -1. 移除回复开头对用户问题的回显（如"你好，你能干吗？ 你好！..."）──────
    # 部分模型会把用户原始问题复述在回复最前面，需要识别并裁剪掉。
    if user_message:
        um = user_message.strip()
        stripped = text.lstrip()
        # 检查回复是否以用户问题开头（可能跟着空格、标点、换行）
        if stripped.lower().startswith(um.lower()):
            after = stripped[len(um):]
            # 跳过紧跟的空白/标点分隔符，保留真正的回复内容
            after = after.lstrip(" \t\r\n，。！？,.!? ")
            text = after

    # ── 0. 移除 AI 模型的思考/推理标签（如 <think>...</think>）──────────────
    # 部分模型（DeepSeek、QwQ 等）会在输出中包含 <think> 标签，需要完整移除。
    text = re.sub(r"(?is)<think>[\s\S]*?</think>", "", text)
    text = re.sub(r"(?is)<thinking>[\s\S]*?</thinking>", "", text)
    text = re.sub(r"(?is)<reasoning>[\s\S]*?</reasoning>", "", text)
    text = re.sub(r"(?is)<reflection>[\s\S]*?</reflection>", "", text)
    # 处理未闭合的思考标签（流式传输中可能出现）
    text = re.sub(r"(?is)<think>[\s\S]*$", "", text)
    text = re.sub(r"(?is)<thinking>[\s\S]*$", "", text)
    text = re.sub(r"(?is)<reasoning>[\s\S]*$", "", text)
    text = re.sub(r"(?is)<reflection>[\s\S]*$", "", text)
    # 移除可能残留的闭合标签
    text = re.sub(r"(?i)</?(think|thinking|reasoning|reflection)>", "", text)

    # ── 1. 清除 codebot 注入的 XML 包装标签（兜底，理论上 system 分离后不再出现）─
    text = re.sub(r"(?is)<(system_policy|conversation_context)>[\s\S]*?</\1>", "", text)
    text = re.sub(r"(?is)</?(system_policy|conversation_context)>", "", text)
    text = re.sub(r"(?is)<internal_context>[\s\S]*?</internal_context>", "", text)
    text = re.sub(r"(?is)<user_message>[\s\S]*?</user_message>", "", text)

    # ── 2. 清除附件注入的上下文块（【附件：...】...```）────────────────────
    text = re.sub(r"(?m)^【附件[：:][^】]*】\n```[\s\S]*?```\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"(?m)^【附件[：:][^】]*（二进制文件，类型：[^）]*）】\s*$", "", text)
    text = re.sub(r"(?m)^【图片附件[：:][^】]*】\n?data:image/[^\s]+;base64,[A-Za-z0-9+/=\r\n]+\s*$", "", text)
    text = re.sub(r"(?m)^【图片附件[：:][^】]*】\s*$", "", text)

    # ── 3. 清除用户消息/输入包装标记行 ─────────────────────────────────────
    text = re.sub(r"(?m)^【用户(?:输入|消息|请求)】.*$", "", text)

    # ── 4. 过滤内部 prompt 标记行（定时任务/提醒/一次性标记）───────────────
    text = re.sub(r"(?m)^__RUN_ONCE__.*$", "", text)
    text = re.sub(r"(?m)^__REMINDER__.*$", "", text)

    # ── 5. 过滤 cron 表达式行（如 "0 9 * * *"）─────────────────────────────
    text = re.sub(r"(?m)^\s*\d+\s+\d+\s+[\d*,/-]+\s+[\d*,/-]+\s+[\d*,/-]+\s*$", "", text)

    text = text.strip()
    return text


def _skill_content_is_noise(text: str) -> bool:
    if not text:
        return True
    lower = text.lower()
    blocked = [
        "system_policy",
        "conversation_context",
        "internal_context",
        "请只输出给用户的最终结果",
    ]
    return any(token in lower for token in blocked)


def _save_chat_log(
    conversation_id: int,
    user_message: str,
    internal_prompt: str,
    tool_events: list,
    final_reply: str,
    model: Optional[str] = None,
    mode: Optional[str] = None,
) -> None:
    """将本次聊天的内部提示词、推理过程和最终回复写入聊天日志表。"""
    try:
        conversations_db.connect()
        cursor = conversations_db.conn.cursor()
        cursor.execute(
            """INSERT INTO chat_logs 
               (conversation_id, user_message, internal_prompt, tool_events, final_reply, model, mode)
               VALUES (?, ?, ?, ?, ?, ?, ?)""",
            (
                conversation_id,
                user_message or "",
                internal_prompt or "",
                json.dumps(tool_events or [], ensure_ascii=False),
                final_reply or "",
                model or "",
                mode or "",
            )
        )
        conversations_db.conn.commit()
    except Exception as e:
        logger.warning(f"保存聊天日志失败: {e}")


def _runtime_state(conv_id: str) -> Dict[str, Any]:
    state = _runtime_stream_state.get(conv_id)
    if state is None:
        state = {"seq": 0, "events": [], "content": "", "running": False, "updated_at": datetime.now()}
        _runtime_stream_state[conv_id] = state
    return state


def _cleanup_runtime_states():
    now = datetime.now()
    to_remove = []
    for conv_id, state in _runtime_stream_state.items():
        updated_at = state.get("updated_at") or now
        running = bool(state.get("running"))
        if not running and (now - updated_at) > timedelta(minutes=20):
            to_remove.append(conv_id)
    for conv_id in to_remove:
        _runtime_stream_state.pop(conv_id, None)


def _runtime_start(conv_id: str):
    _cleanup_runtime_states()
    state = _runtime_state(conv_id)
    state["seq"] = 0
    state["events"] = []
    state["content"] = ""
    state["running"] = True
    state["updated_at"] = datetime.now()


def _runtime_set_content(conv_id: str, content: str):
    state = _runtime_state(conv_id)
    state["content"] = content or ""
    state["updated_at"] = datetime.now()


def _runtime_append_event(conv_id: str, event: dict):
    state = _runtime_state(conv_id)
    state["seq"] = int(state.get("seq", 0)) + 1
    payload = dict(event)
    payload["seq"] = state["seq"]
    payload["created_at"] = datetime.utcnow().isoformat() + "Z"
    events = state.get("events", [])
    events.append(payload)
    if len(events) > 1000:
        events = events[-1000:]
    state["events"] = events
    state["updated_at"] = datetime.now()


def _runtime_finish(conv_id: str, content: str = ""):
    state = _runtime_state(conv_id)
    if content:
        state["content"] = content
    state["running"] = False
    state["updated_at"] = datetime.now()


def _runtime_snapshot(conv_id: str, since_seq: int = 0) -> Dict[str, Any]:
    state = _runtime_stream_state.get(conv_id)
    if not state:
        return {"events": [], "last_seq": int(since_seq), "content": "", "running": False}
    events = state.get("events", [])
    filtered = [e for e in events if int(e.get("seq", 0)) > int(since_seq)]
    return {
        "events": filtered,
        "last_seq": int(state.get("seq", 0)),
        "content": state.get("content", "") or "",
        "running": bool(state.get("running")),
    }


def _is_multi_agent_hub(conversation: Optional[Dict]) -> bool:
    return bool(conversation and conversation.get("conversation_type") == "multi_agent_hub")


def _conversation_role_label(conversation: Dict) -> str:
    role = str(conversation.get("group_role") or "").strip()
    if role:
        return role
    title = str(conversation.get("title") or "").strip()
    return title or f"Agent-{conversation.get('id')}"


def _find_member_by_label(label: str, members: List[Dict]) -> Optional[Dict]:
    target = (label or "").strip().lower()
    if not target:
        return None
    for member in members:
        role = str(member.get("group_role") or "").strip().lower()
        title = str(member.get("title") or "").strip().lower()
        if target in {role, title} or target in role or target in title:
            return member
    return None


def _pick_member_for_task(task_text: str, members: List[Dict]) -> Dict:
    text = (task_text or "").lower()
    explicit_targets = re.findall(r"(?:分配给|交给|给)\s*[“\"']([^”\"']+)[”\"']\s*(?:对话|agent|Agent)?", task_text or "")
    for target in explicit_targets:
        member = _find_member_by_label(target, members)
        if member:
            return member
    keyword_map = [
        ("鉴赏", ["鉴赏", "赏析", "点评", "评价", "review", "critique"]),
        ("评论", ["评论", "点评", "评价", "review", "critique"]),
        ("写作", ["写", "创作", "文案", "文章", "诗", "故事", "小说", "write"]),
        ("前端", ["前端", "vue", "react", "ui", "页面", "组件", "样式", "css", "frontend"]),
        ("后端", ["后端", "api", "接口", "fastapi", "python", "服务", "backend"]),
        ("数据库", ["数据库", "sqlite", "sql", "表", "字段", "迁移", "database", "db"]),
        ("测试", ["测试", "验证", "用例", "test", "pytest", "build", "检查"]),
    ]
    for preferred_role, words in keyword_map:
        if not any(word in text for word in words):
            continue
        for member in members:
            haystack = f"{member.get('title') or ''} {member.get('group_role') or ''}".lower()
            if preferred_role in haystack or any(word in haystack for word in words):
                return member
    return members[0]


def _split_multi_agent_task_text(user_message: str) -> List[str]:
    lines = [line.strip("- 　\t") for line in re.split(r"[\n；;]+", user_message or "") if line.strip()]
    if len(lines) <= 1:
        sentences = [s.strip() for s in re.split(r"(?<=[。！？!?])", user_message or "") if s.strip()]
        lines = sentences if len(sentences) > 1 else [user_message.strip()]
    return [line for line in lines if line]


def _extract_explicit_multi_agent_steps(user_message: str, members: List[Dict]) -> List[Dict[str, Any]]:
    text = user_message or ""
    pattern = re.compile(
        r"(?:将|把)\s*(?P<task>.+?)\s*(?:分配给|交给)\s*[“\"'](?P<label>[^”\"']+)[”\"']\s*(?:对话|agent|Agent)?"
    )
    matches = list(pattern.finditer(text))
    assignments: List[Dict[str, Any]] = []
    for idx, match in enumerate(matches):
        label = match.group("label") or ""
        member = _find_member_by_label(label, members)
        if not member:
            continue
        task_text = (match.group("task") or "").strip(" ，,。；;\n\t")
        prefix = text[:match.start()].strip(" ，,。；;\n\t") if idx == 0 else ""
        next_start = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        tail = text[match.end():next_start].strip(" ，,。；;\n\t")
        tail = re.sub(r"^(?:写完.*?之后|写完.*?后|完成.*?之后|完成.*?后|然后|再|接着|随后|之后|最后)", "", tail).strip(" ，,。；;\n\t")
        if prefix and task_text and len(prefix) <= 80:
            task_text = f"{prefix}；明确负责：{task_text}"
        if tail:
            task_text = f"{task_text}；{tail}" if task_text else tail
        if not task_text:
            task_text = f"按用户要求处理分配给 {label} 的任务"
        assignments.append({"task_id": len(assignments) + 1, "task": task_text, "member": member})
    return assignments


def _plan_multi_agent_tasks(user_message: str, members: List[Dict]) -> List[List[Dict[str, Any]]]:
    text = (user_message or "").strip()
    if not text:
        return []

    explicit_assignments = _extract_explicit_multi_agent_steps(text, members)
    has_sequential_signal = bool(re.search(r"写完|完成后|完成之后|然后|再|接着|随后|之后|最后|交给", text))
    if explicit_assignments:
        if has_sequential_signal:
            return [[{**assignment, "depends_on_previous": idx > 0}] for idx, assignment in enumerate(explicit_assignments)]
        return [[{**assignment, "depends_on_previous": False} for assignment in explicit_assignments]]

    # Chinese prompts often encode dependencies with "then/after" in one sentence.
    sequential_parts = [p.strip(" ，,。；;\n\t") for p in re.split(r"(?:写完(?:之后|后)?|完成(?:之后|后)?|然后|再|接着|随后|之后|最后|交给)", text) if p.strip(" ，,。；;\n\t")]
    raw_tasks = sequential_parts if has_sequential_signal and len(sequential_parts) > 1 else _split_multi_agent_task_text(text)

    steps: List[List[Dict[str, Any]]] = []
    if has_sequential_signal and len(raw_tasks) > 1:
        for idx, task_text in enumerate(raw_tasks, 1):
            member = _pick_member_for_task(task_text, members)
            steps.append([{"task_id": idx, "task": task_text, "member": member, "depends_on_previous": idx > 1}])
        return steps

    parallel_step: List[Dict[str, Any]] = []
    for idx, task_text in enumerate(raw_tasks, 1):
        member = _pick_member_for_task(task_text, members)
        parallel_step.append({"task_id": idx, "task": task_text, "member": member, "depends_on_previous": False})
    return [parallel_step] if parallel_step else []


def _split_multi_agent_tasks(user_message: str, members: List[Dict]) -> List[Dict[str, Any]]:
    return [assignment for step in _plan_multi_agent_tasks(user_message, members) for assignment in step]


def _format_multi_agent_plan(user_message: str, steps: List[List[Dict[str, Any]]]) -> str:
    lines = ["## 多Agent任务计划", "", f"用户任务：{user_message}", ""]
    if not steps:
        lines.append("暂无可执行任务。")
        return "\n".join(lines).strip()
    for step_index, step in enumerate(steps, 1):
        mode = "并行" if len(step) > 1 else "串行"
        lines.append(f"### 第 {step_index} 步（{mode}）")
        for assignment in step:
            member = assignment["member"]
            role = _conversation_role_label(member)
            dependency = "，依赖上一轮产物" if assignment.get("depends_on_previous") else ""
            lines.append(f"- {role}（对话 #{member.get('id')}）：{assignment['task']}{dependency}")
        lines.append("")
    return "\n".join(lines).strip()


def _append_hub_progress(conv_id: str, lines: List[str], message: str, event_type: str = "status", **extra):
    lines.append(message)
    content = "\n".join(lines).strip()
    _runtime_set_content(conv_id, content)
    payload = {"type": event_type, "message": message, **extra}
    _runtime_append_event(conv_id, payload)


async def _build_multi_agent_hub_reply(user_message: str, members: List[Dict], results: List[Dict[str, str]], plan_text: str = "") -> str:
    lines = [
        "## 多Agent群聊任务结果",
        "",
        f"用户任务：{user_message}",
        "",
    ]
    if plan_text:
        lines.extend([plan_text, ""])
    lines.append("### 执行过程")
    for item in results:
        lines.append(f"- 第 {item.get('step', 1)} 步：{item['role']}（对话 #{item['conversation_id']}）完成 `{item['task']}`")
    lines.extend(["", "### 成员回复"])
    for item in results:
        lines.append(f"#### {item['role']}")
        lines.append(item.get("reply") or "未返回内容")
        lines.append("")
    member_names = "、".join(_conversation_role_label(m) for m in members) or "暂无成员"
    lines.extend(["### 协作状态", f"参与成员：{member_names}"])
    return "\n".join(lines).strip()


async def _execute_opencode_client(message: str, model: Optional[str] = None, mode: Optional[str] = None, conversation_id: Optional[str] = None, system: Optional[str] = None, user_message: str = "") -> Tuple[Optional[str], bool]:
    client = opencode_ws
    created_client = False
    try:
        if client is None:
            client = OpenCodeClient(app_config.opencode.server_url)
            created_client = True
        ok = await _ensure_opencode_client_connected(client)
        if not ok:
            return None, False
        result = await client.execute_task(message, model=model, mode=mode, conversation_id=conversation_id, system=system)
        if result.success:
            return _sanitize_assistant_output(result.content or "", user_message=user_message) or None, True
        return result.error or None, True
    except Exception as e:
        logger.error(f"OpenCode 调用失败: {e}")
        return None, False
    finally:
        if created_client and client and getattr(client, "connected", False):
            try:
                await client.disconnect()
            except Exception:
                pass


async def _execute_opencode_client_with_parts(
    message: str,
    model: Optional[str] = None,
    mode: Optional[str] = None,
    conversation_id: Optional[str] = None,
    system: Optional[str] = None,
    user_message: str = ""
) -> Tuple[Optional[str], bool, List[dict]]:
    client = opencode_ws
    created_client = False
    try:
        if client is None:
            client = OpenCodeClient(app_config.opencode.server_url)
            created_client = True
        ok = await _ensure_opencode_client_connected(client)
        if not ok:
            return None, False, []
        result = await client.execute_task(message, model=model, mode=mode, conversation_id=conversation_id, system=system)
        if result.success:
            return _sanitize_assistant_output(result.content or "", user_message=user_message) or None, True, result.parts or []
        return result.error or None, True, []
    except Exception as e:
        logger.error(f"OpenCode 调用失败: {e}")
        return None, False, []
    finally:
        if created_client and client and getattr(client, "connected", False):
            try:
                await client.disconnect()
            except Exception:
                pass


# ── MCP 聊天意图识别与处理 ─────────────────────────────────────────────────

def _looks_like_mcp_message(message: str) -> bool:
    """
    判断消息是否与 MCP Server **管理**相关（添加/删除/启用/禁用等）。

    规则：
    1. 必须同时含有 MCP 关键词 AND 管理动作关键词
    2. 排除"调用/使用/invoke/call/run/执行"等使用意图——这类消息应透传给 opencode 处理
    """
    if not message:
        return False
    text = message.strip()
    lower = text.lower()

    mcp_keywords = ["mcp", "model context protocol"]
    has_mcp = any(kw in lower for kw in mcp_keywords)
    if not has_mcp:
        return False

    # 明确的"使用/调用"意图 → 不拦截，让 opencode 处理
    invoke_keywords = ["调用", "使用", "invoke", "call", "run", "执行", "帮我用", "通过mcp", "通过 mcp",
                       "用mcp", "用 mcp", "利用", "借助"]
    if any(kw in lower for kw in invoke_keywords):
        return False

    # 必须同时有管理动作才算管理消息
    mgmt_action_keywords = ["添加", "新增", "安装", "配置", "删除", "移除", "卸载", "启用", "禁用",
                             "列出", "查看", "查询", "add", "install", "remove", "delete",
                             "enable", "disable", "list"]
    has_mgmt = any(kw in lower for kw in mgmt_action_keywords)
    if not has_mgmt:
        return False

    # 还需要有"服务器/server"语境，避免误杀含 mcp 词的普通句子
    has_server_ctx = bool(re.search(r"(server|服务器|服务|mcp\s*服务)", lower))
    return has_server_ctx


def _extract_mcp_command_from_message(message: str) -> Optional[dict]:
    """
    尝试从聊天消息中提取 MCP Server 配置。
    支持格式示例：
      - 添加 MCP 服务器 名称=filesystem 命令=npx 参数=-y @modelcontextprotocol/server-filesystem
      - 帮我配置一个 MCP，命令：uvx mcp-server-git
      - 新增 SSE 类型 MCP 服务器，URL=http://localhost:3000/sse
    返回 dict 或 None（解析失败）。
    """
    text = message.strip()

    # 优先尝试 JSON 格式（用户直接粘贴）
    json_match = re.search(r"\{[\s\S]*\}", text)
    if json_match:
        try:
            data = json.loads(json_match.group(0))
            if isinstance(data, dict) and ("command" in data or "url" in data):
                return {
                    "name": str(data.get("name") or data.get("id") or "").strip(),
                    "description": str(data.get("description") or "").strip(),
                    "transport": str(data.get("transport") or "stdio"),
                    "command": str(data.get("command") or "").strip() or None,
                    "args": data.get("args") or [],
                    "url": str(data.get("url") or "").strip() or None,
                    "env": data.get("env") or {},
                }
        except Exception:
            pass

    result: dict = {
        "name": "",
        "description": "",
        "transport": "stdio",
        "command": None,
        "args": [],
        "url": None,
        "env": {},
    }

    # SSE URL 识别
    url_match = re.search(r"url\s*[=:=]\s*(\S+)", text, re.IGNORECASE)
    if not url_match:
        url_match = re.search(r"https?://\S+", text)
    if url_match:
        result["url"] = url_match.group(1) if "=" in url_match.group(0) else url_match.group(0)
        result["transport"] = "sse"

    # stdio 命令识别
    cmd_match = re.search(
        r"(?:命令|command)\s*[=:＝]\s*([^\n，,;；]+)",
        text, re.IGNORECASE
    )
    if cmd_match:
        cmd_str = cmd_match.group(1).strip()
        parts = cmd_str.split()
        result["command"] = parts[0] if parts else None
        result["args"] = parts[1:] if len(parts) > 1 else []
    elif not result["url"]:
        # 尝试识别常见命令行模式：npx/uvx/node/python
        cmd_inline = re.search(
            r"\b(npx|uvx|node|python|python3|uv|deno|bun)\s+([^\n，,;；\"\']+)",
            text
        )
        if cmd_inline:
            parts = cmd_inline.group(0).split()
            result["command"] = parts[0]
            result["args"] = parts[1:]

    # 名称识别
    name_match = re.search(r"(?:名称|name|叫做|叫)\s*[=:＝]?\s*([^\s，,;；\n]+)", text, re.IGNORECASE)
    if name_match:
        result["name"] = name_match.group(1).strip()
    elif result["command"]:
        # 从命令自动命名
        cmd_base = result["args"][-1] if result["args"] else result["command"]
        result["name"] = re.sub(r"[@/].*", "", cmd_base).strip() or result["command"]
    elif result["url"]:
        try:
            from urllib.parse import urlparse as _urlparse
            result["name"] = _urlparse(result["url"]).hostname or "mcp-sse"
        except Exception:
            result["name"] = "mcp-sse"

    # 描述识别
    desc_match = re.search(r"(?:描述|description|说明)\s*[=:＝]\s*([^\n，,;；]+)", text, re.IGNORECASE)
    if desc_match:
        result["description"] = desc_match.group(1).strip()

    if not result["command"] and not result["url"]:
        return None

    return result


async def _try_handle_mcp_message(message: str) -> Optional[str]:
    """尝试将聊天消息路由为 MCP 操作，成功返回回复字符串，否则返回 None。"""
    if not _looks_like_mcp_message(message):
        return None

    text = message.strip()
    lower = text.lower()

    # ── 查询/列出 ──────────────────────────────────────────────────────────
    list_keywords = ["列出", "查看", "查询", "显示", "有哪些", "有什么", "list"]
    if any(kw in lower for kw in list_keywords):
        from api.routes.mcp import _read_all
        servers = _read_all()
        if not servers:
            return "当前没有配置任何 MCP Server。\n可在「MCP」页面手动添加，也可以告诉我想添加什么 MCP Server。"
        lines = [f"当前共有 {len(servers)} 个 MCP Server："]
        for s in servers:
            state = "启用" if s.get("enabled", True) else "禁用"
            transport = s.get("transport", "stdio")
            cmd_or_url = s.get("command") or s.get("url") or ""
            lines.append(f"- [{state}] **{s['name']}** ({transport}) `{cmd_or_url}`")
        lines.append("\n可在「MCP」页面管理。")
        return "\n".join(lines)

    # ── 删除/移除 ──────────────────────────────────────────────────────────
    # 必须同时有删除关键词 AND MCP Server 语境，防止误拦截"修改文件内容"等正常请求
    delete_keywords = ["删除", "移除", "卸载", "remove", "delete"]
    has_server_ctx_del = bool(re.search(r"(mcp\s*(server|服务器)|server.*mcp|服务器.*mcp)", lower))
    if any(kw in lower for kw in delete_keywords) and has_server_ctx_del:
        from api.routes.mcp import _read_all, _write_all
        servers = _read_all()
        # 尝试按名称匹配
        name_match = re.search(
            r"(?:\u5220\u9664|\u79fb\u9664|\u5378\u8f7d|remove|delete)\s+(?:mcp\s*)?(?:\u670d\u52a1\u5668?|server)?\s*\S*?([A-Za-z0-9_\-\.]+)",
            text, re.IGNORECASE
        )
        if not name_match:
            # 更宽松匹配：删除/移除 后面的第一个非空词
            name_match = re.search(
                r"(?:\u5220\u9664|\u79fb\u9664|\u5378\u8f7d|remove|delete)[^\w]*([A-Za-z0-9_\-\.]+)",
                text, re.IGNORECASE
            )
        if name_match:
            target_name = name_match.group(1).strip().strip("\"'")
            matched = [s for s in servers if target_name.lower() in s["name"].lower()]
            if matched:
                s = matched[0]
                servers = [sv for sv in servers if sv["id"] != s["id"]]
                _write_all(servers)
                return f"已删除 MCP Server「{s['name']}」。"
            return f"未找到名称包含「{target_name}」的 MCP Server，请检查名称或在「MCP」页面操作。"
        return "请告诉我要删除哪个 MCP Server 的名称，例如：删除 MCP 服务器 filesystem"

    # ── 添加/新增/安装/配置 ────────────────────────────────────────────────
    add_keywords = ["添加", "新增", "安装", "配置", "add", "install", "register", "注册"]
    if any(kw in lower for kw in add_keywords) or re.search(
        r"(npx|uvx|node|python)\s+", text
    ):
        parsed = _extract_mcp_command_from_message(message)
        if not parsed:
            return (
                "我识别到你想添加 MCP Server，但没能解析出有效的命令或 URL。\n"
                "请参考以下格式：\n"
                "- **stdio 模式**：`添加 MCP 服务器，命令=npx -y @modelcontextprotocol/server-filesystem`\n"
                "- **SSE 模式**：`添加 MCP 服务器，URL=http://localhost:3000/sse`\n"
                "也可以在「MCP」页面手动填写。"
            )
        try:
            server = mcp_router.add_server_from_chat(
                name=parsed["name"] or "未命名",
                transport=parsed["transport"],
                command=parsed["command"],
                args=parsed["args"],
                url=parsed["url"],
                env=parsed["env"],
                description=parsed["description"],
            )
            transport_label = "SSE" if parsed["transport"] == "sse" else "stdio"
            detail = f"URL：`{parsed['url']}`" if parsed["url"] else f"命令：`{parsed['command']} {' '.join(parsed['args'])}`.strip()"
            return (
                f"已添加 MCP Server「{server['name']}」（{transport_label} 模式）。\n"
                f"{detail}\n"
                "可在「MCP」页面查看和管理。"
            )
        except Exception as e:
            logger.error(f"聊天添加 MCP Server 失败: {e}")
            return f"添加 MCP Server 失败：{e}\n请在「MCP」页面手动添加。"

    return None


async def _execute_sandbox(prompt: str) -> Optional[str]:
    """尝试在沙箱 VM 中执行任务，返回结果字符串；沙箱未就绪时返回 None（降级到本地执行）。"""
    if sandbox_manager is None:
        return None
    if not app_config.sandbox.enabled:
        return None
    try:
        result = await sandbox_manager.execute(prompt)
        if result.success:
            return result.content or ""
        logger.warning(f"沙箱执行失败（降级本地）: {result.error}")
        return None
    except Exception as e:
        logger.warning(f"沙箱执行异常（降级本地）: {e}")
        return None


def _should_use_sandbox(message: str) -> bool:
    """根据执行模式配置决定是否路由到沙箱。"""
    mode = getattr(app_config.sandbox, "execution_mode", "auto")
    if not app_config.sandbox.enabled:
        return False
    if mode == "sandbox":
        return True
    if mode == "local":
        return False
    # auto 模式：含文件操作/代码执行等高风险关键词时走沙箱
    if mode == "auto":
        risky_patterns = [
            r"执行代码", r"运行脚本", r"execute\s+code", r"run\s+script",
            r"rm\s+-rf", r"删除文件", r"格式化", r"sudo", r"chmod",
            r"pip\s+install", r"npm\s+install", r"apt[-\s]",
        ]
        import re as _re
        for pat in risky_patterns:
            if _re.search(pat, message, _re.IGNORECASE):
                return True
    return False


async def _sync_codebot_as_third_party():
    mcp_router.auto_sync_mcp_to_opencode()
    mcp_router.ensure_codebot_remote_mcp_in_opencode()


async def _execute_hermes_proxy(
    message: str,
    model: Optional[str] = None,
    conversation_id: Optional[str] = None,
    knowledge_paths: Optional[List[str]] = None,
) -> str:
    from api.routes import hermes as hermes_router

    selected_skill = None
    cleaned_message = message
    hermes_skills: List[str] = []
    try:
        selected_skill, cleaned_message, _ = _extract_requested_skill(message)
        if selected_skill:
            skill_name = selected_skill.get("slug") or selected_skill.get("name") or ""
            if skill_name:
                hermes_skills.append(str(skill_name))
    except Exception as exc:
        logger.debug(f"[Hermes] skill marker parse failed: {exc}")

    hermes_message = cleaned_message or message
    if knowledge_paths:
        obsidian_context = _build_obsidian_context(hermes_message, knowledge_paths)
        if obsidian_context:
            hermes_message = (
                f"{hermes_message}\n\n"
                "[Codebot selected Obsidian Markdown context]\n"
                f"{obsidian_context}"
            )
            if "obsidian" not in {name.lower() for name in hermes_skills}:
                hermes_skills.append("obsidian")

    if _looks_like_codebot_schedule_creation_request(hermes_message):
        hermes_message = (
            f"{_build_codebot_scheduler_boundary('hermes')}\n\n"
            "User request:\n"
            f"{hermes_message}"
        )

    conv_id = str(conversation_id) if conversation_id is not None else ""
    if conv_id:
        mark_conversation_running(conv_id)
    try:
        response = await hermes_router.hermes_chat(
            hermes_router.HermesChatRequest(
                message=hermes_message,
                model=model,
                conversation_id=conversation_id,
                skills=hermes_skills,
            )
        )
        return ((response.get("data") or {}).get("content") or "").strip()
    finally:
        if conv_id:
            unmark_conversation_running(conv_id)


async def _execute_opencode(
    message: str,
    model: Optional[str] = None,
    mode: Optional[str] = None,
    conversation_id: Optional[str] = None,
    project_dir: Optional[str] = None,
    target: Optional[str] = None,
    knowledge_paths: Optional[List[str]] = None,
) -> str:
    schedule_result = await _handle_codebot_schedule_creation_request(
        message,
        target=target,
        conversation_id=int(conversation_id) if str(conversation_id or "").isdigit() else None,
        execution_model=model,
    )
    if schedule_result:
        return schedule_result

    if (target or "").strip().lower() == "hermes":
        content = await _execute_hermes_proxy(message, model=model, conversation_id=conversation_id, knowledge_paths=knowledge_paths)
        return _sanitize_assistant_output(content or "", user_message=message)

    try:
        await _sync_codebot_as_third_party()
    except Exception as _sync_err:
        logger.debug(f"[Codebot] 第三方能力同步失败（跳过）: {_sync_err}")

    skill_watch = _begin_codebot_skill_creation_watch(message)
    # 系统指令通过独立 system 字段传递，用户消息保持纯净
    system_prompt, user_message = await _build_opencode_prompt_parts(
        message,
        mode=mode,
        project_dir=project_dir,
        target=target,
        knowledge_paths=knowledge_paths,
        model=model,
    )
    try:
        content, _ = await _execute_opencode_client(user_message, model=model, mode=mode, conversation_id=conversation_id, system=system_prompt, user_message=message)
        content = _sanitize_assistant_output(content or "", user_message=message)

        if not content:
            content = "OpenCode 未连接，请先启动本地 opencode server 或检查 server_url 配置"

        return content
    finally:
        _finish_codebot_skill_creation_watch(skill_watch)


async def _execute_opencode_with_meta(
    message: str,
    model: Optional[str] = None,
    mode: Optional[str] = None,
    conversation_id: Optional[str] = None,
    project_dir: Optional[str] = None,
    target: Optional[str] = None,
    knowledge_paths: Optional[List[str]] = None,
) -> Tuple[str, List[dict]]:
    schedule_result = await _handle_codebot_schedule_creation_request(
        message,
        target=target,
        conversation_id=int(conversation_id) if str(conversation_id or "").isdigit() else None,
        execution_model=model,
    )
    if schedule_result:
        return schedule_result, []

    if (target or "").strip().lower() == "hermes":
        content = await _execute_hermes_proxy(message, model=model, conversation_id=conversation_id, knowledge_paths=knowledge_paths)
        return _sanitize_assistant_output(content or "", user_message=message), []

    try:
        await _sync_codebot_as_third_party()
    except Exception as _sync_err:
        logger.debug(f"[Codebot] 第三方能力同步失败（跳过）: {_sync_err}")

    skill_watch = _begin_codebot_skill_creation_watch(message)
    # 系统指令通过独立 system 字段传递，用户消息保持纯净
    system_prompt, user_message = await _build_opencode_prompt_parts(
        message,
        mode=mode,
        project_dir=project_dir,
        target=target,
        knowledge_paths=knowledge_paths,
        model=model,
    )
    try:
        content, _, parts = await _execute_opencode_client_with_parts(
            user_message,
            model=model,
            mode=mode,
            conversation_id=conversation_id,
            system=system_prompt,
            user_message=message
        )
        content = _sanitize_assistant_output(content or "", user_message=message)

        if not content:
            content = "OpenCode 未连接，请先启动本地 opencode server 或检查 server_url 配置"
            return content, []

        return content, parts or []
    finally:
        _finish_codebot_skill_creation_watch(skill_watch)


async def _stream_execute_opencode_with_meta(
    message: str,
    model: Optional[str] = None,
    mode: Optional[str] = None,
    conversation_id: Optional[str] = None,
    project_dir: Optional[str] = None,
    target: Optional[str] = None,
    knowledge_paths: Optional[List[str]] = None,
):
    if _looks_like_codebot_schedule_creation_request(message):
        executor = _task_executor_from_target(target)
        preparing = f"Codebot 正在创建定时任务（执行器：{_task_executor_label(executor)}）...\n\n"
        yield {
            "type": "content_delta",
            "delta": preparing,
            "content": preparing,
            "source": "codebot_scheduler",
        }
        content = await _handle_codebot_schedule_creation_request(
            message,
            target=target,
            conversation_id=int(conversation_id) if str(conversation_id or "").isdigit() else None,
            execution_model=model,
        )
        content = content or "Codebot 未能创建该定时任务，请补充明确时间后重试。"
        yield {
            "type": "done",
            "content": content,
            "parts": [],
            "source": "codebot_scheduler",
            "agent": "Codebot Scheduler",
        }
        return

    if (target or "").strip().lower() == "hermes":
        model_text = f"（模型：{model}）" if model else ""
        yield {
            "type": "content_delta",
            "delta": f"Hermes Agent CLI 正在处理{model_text}...\n\n",
            "content": f"Hermes Agent CLI 正在处理{model_text}...\n\n",
            "source": "hermes",
        }
        content = await _execute_hermes_proxy(message, model=model, conversation_id=conversation_id, knowledge_paths=knowledge_paths)
        yield {
            "type": "status",
            "phase": "hermes_cli",
            "source": "hermes",
            "message": "Hermes Agent CLI completed",
        }
        yield {
            "type": "done",
            "content": _sanitize_assistant_output(content or "", user_message=message),
            "parts": [],
            "source": "hermes",
            "agent": "Hermes Agent CLI",
        }
        return

    try:
        await _sync_codebot_as_third_party()
    except Exception as _sync_err:
        logger.debug(f"[Codebot] 第三方能力同步失败（跳过）: {_sync_err}")

    skill_watch = _begin_codebot_skill_creation_watch(message)
    # 系统指令通过独立 system 字段传递，用户消息保持纯净
    system_prompt, user_message = await _build_opencode_prompt_parts(
        message,
        mode=mode,
        project_dir=project_dir,
        target=target,
        knowledge_paths=knowledge_paths,
        model=model,
    )
    # 首先 yield 内部提示词事件，供聊天日志记录
    yield {"type": "internal_prompt", "prompt": f"[system]\n{system_prompt}\n\n[user]\n{user_message}"}

    client = opencode_ws
    created_client = False
    try:
        if client is None:
            client = OpenCodeClient(app_config.opencode.server_url)
            created_client = True
        ok = await _ensure_opencode_client_connected(client)
        if not ok:
            yield {"type": "done", "content": "OpenCode 未连接，请先启动本地 opencode server 或检查 server_url 配置", "parts": []}
            return
        async for event in client.execute_task_stream(
            prompt=user_message,
            model=model,
            mode=mode,
            conversation_id=conversation_id,
            system=system_prompt
        ):
            yield event
    finally:
        _finish_codebot_skill_creation_watch(skill_watch)
        if created_client and client and getattr(client, "connected", False):
            try:
                await client.disconnect()
            except Exception:
                pass


def _extract_json_object(text: str) -> Optional[dict]:
    if not text:
        return None
    match = re.search(r"\{[\s\S]*\}", text)
    if not match:
        return None
    candidate = match.group(0)
    try:
        return json.loads(candidate)
    except Exception:
        return None


async def _try_ai_route_action(message: str) -> Tuple[Optional[str], bool]:
    opencode_available = False
    if not message or not message.strip():
        return None, opencode_available

    memory_like = _looks_like_memory_message(message)
    schedule_like = _looks_like_schedule_message(message)
    birthday_reminder_like = _looks_like_birthday_reminder_message(message)
    # 定时任务优先：已识别为定时任务则不再判断记忆
    if schedule_like and not birthday_reminder_like:
        memory_like = False
    allowed_actions = []
    if schedule_like:
        allowed_actions.append("create_scheduled_task")
    if memory_like:
        allowed_actions.append("save_memory")
    if not allowed_actions:
        return None, opencode_available

    prompt = (
        "你是意图识别与结构化指令提取器，用于辅助 OpenCode 桌面工作流。\n"
        "任务：从用户输入中提取结构化指令。\n"
        f"本次允许的 action 只有：{ '|'.join(allowed_actions) }。\n"
        "只输出 JSON（不要输出任何解释、不要输出 Markdown、不要输出代码块）。\n"
        "输出格式：\n"
        "{"
        "\"action\":\"create_scheduled_task\"|\"save_memory\","
        "\"confidence\":0.0,"
        "\"task\":{"
        "\"name\":\"\","
        "\"cron\":\"\","
        "\"cron_prompt\":\"\","
        "\"task_prompt\":\"\","
        "\"run_once\":false,"
        "\"notify_channels\":[\"app\"]"
        "},"
        "\"memory\":{"
        "\"content\":\"\","
        "\"category\":\"note\""
        "},"
        "\"reply\":\"\""
        "}\n"
        "要求：\n"
        "1) 当允许 action 里包含 create_scheduled_task 且用户表达“到某个时间/周期要做某事”（例如提醒/生成/写作/整理/保存）并包含明确时间（例如 HH:MM 或 X点），必须选择 action=create_scheduled_task。\n"
        "2) action=create_scheduled_task 时，task.task_prompt 必须是未来执行时要做的事（例如“提醒用户：xxx”），不要在 reply 中输出任何自然语言回复。\n"
        "3) 如果用户表达“今天/明天/后天/只提醒一次/仅一次”，则 run_once=true。\n"
        "4) task.cron 是标准 5 段 cron（分 时 日 月 周），能给则给；不确定就把时间表达写入 task.cron_prompt。\n"
        "5) action=save_memory 时，将要记住的核心内容放入 memory.content，不要生成多余回复。"
        "memory.category 必须准确选择：preference=喜好/偏爱/风格/工具偏好；habit=日常习惯/行为模式；"
        "profile=仅限姓名/年龄/生日/职业等身份信息；contact=仅限电话/邮箱/微信等联系方式；"
        "address=仅限物理地址；note=其他备忘。不要把偏好或习惯归为profile或contact！\n"
        "6) 只允许输出 JSON，对 JSON 以外的任何字符都视为失败。\n"
        f"用户输入：{message}"
    )

    result_content, opencode_available = await _execute_opencode_client(prompt)
    if not result_content:
        return None, opencode_available
    data = _extract_json_object(result_content)
    if not isinstance(data, dict):
        repair_prompt = (
            "你的上一条输出不是合法 JSON。"
            "请仅输出一个 JSON 对象，必须以 { 开头、以 } 结尾，且不包含任何额外字符。\n"
            f"允许的 action 只有：{ '|'.join(allowed_actions) }。\n"
            "JSON 格式同上一次要求。\n"
            f"用户输入：{message}\n"
            f"你的上一条输出：{result_content}"
        )
        repaired_content, opencode_available = await _execute_opencode_client(repair_prompt)
        if not repaired_content:
            return None, opencode_available
        data = _extract_json_object(repaired_content)
        if not isinstance(data, dict):
            return None, opencode_available

    action = str(data.get("action") or "").strip()
    confidence = data.get("confidence")
    try:
        confidence_value = float(confidence)
    except Exception:
        confidence_value = 0.0

    if confidence_value < 0.35 and len(allowed_actions) > 1:
        return None, opencode_available

    if action not in allowed_actions:
        return None, opencode_available

    if action == "create_scheduled_task":
        if not scheduler_router.scheduler:
            return None, opencode_available
        task_data = data.get("task") if isinstance(data.get("task"), dict) else {}
        name = str(task_data.get("name") or "").strip()
        cron_expression = str(task_data.get("cron") or "").strip()
        cron_prompt = str(task_data.get("cron_prompt") or "").strip()
        task_prompt = str(task_data.get("task_prompt") or "").strip()
        notify_channels = task_data.get("notify_channels")
        if not isinstance(notify_channels, list) or not notify_channels:
            notify_channels = ["app"]
        run_once = (
            bool(task_data.get("run_once"))
            or any(key in (message or "") for key in ["今天", "明天", "后天", "一次", "只提醒一次", "仅提醒一次", "提醒一次"])
            or bool(re.search(r"\d+\s*(分钟|小时|天)\s*(后|之后|以后)", message or ""))
            or bool(re.search(r"(半小时|一小时|一天)\s*(后|之后|以后)", message or ""))
        )

        if not task_prompt:
            content = _extract_reminder_content(message)
            task_prompt = f"请提醒用户：{content}" if content else message

        if run_once:
            task_prompt = f"__RUN_ONCE__\n{task_prompt}"

        if not cron_expression:
            cron_source = cron_prompt or message
            cron_payload = await scheduler_router.ai_generate_cron(prompt=cron_source)
            cron_data = cron_payload.get("data") if isinstance(cron_payload, dict) else None
            cron_expression = cron_data.get("cron") if cron_data else None
        if not cron_expression:
            return None, opencode_available

        if not name:
            content = _extract_reminder_content(message)
            name = f"{'一次性' if run_once else ''}提醒：{content}" if content else ("一次性定时提醒" if run_once else "定时提醒")
        if len(name) > 30:
            name = f"{name[:30]}..."

        task = scheduler_router.scheduler.create_task(
            name=name,
            cron_expression=cron_expression,
            task_prompt=task_prompt,
            notify_channels=notify_channels,
            executor="opencode",
            execution_model=_task_execution_model_from_chat_model(None),
        )
        next_run = task.next_run.isoformat() if task.next_run else "待计算"
        return f"已创建定时任务：{task.name}\nCron：{task.cron_expression}\n下次运行：{next_run}\n可在“定时任务”查看和管理。", opencode_available

    if action == "save_memory":
        mem = data.get("memory") if isinstance(data.get("memory"), dict) else {}
        content = str(mem.get("content") or "").strip()
        category = str(mem.get("category") or "note").strip() or "note"
        if not content:
            return None, opencode_available
        await _save_memory_content(content=content, category=category, raw_message=message)
        return f"好的，我已经记住了：{content}\n可在“记忆”页面查看。", opencode_available

    return None, opencode_available


def _looks_like_birthday_reminder_message(message: str) -> bool:
    text = (message or "").strip()
    if not text:
        return False
    if "生日" not in text:
        return False
    if not _extract_birthday_value(text):
        return False
    reminder_keys = ["提醒", "闹钟", "提示我", "叫我", "通知我", "别忘了", "记得"]
    return any(key in text for key in reminder_keys)


def _extract_time_for_reminder(message: str, default_hour: int = 9, default_minute: int = 0) -> Tuple[int, int]:
    text = (message or "").strip()
    if not text:
        return default_hour, default_minute
    time_match = re.search(r"(\d{1,2})\s*[:：]\s*(\d{1,2})", text)
    if time_match:
        hour = max(0, min(23, int(time_match.group(1))))
        minute = max(0, min(59, int(time_match.group(2))))
        return hour, minute
    half_match = re.search(r"(\d{1,2})\s*点\s*半", text)
    if half_match:
        hour = int(half_match.group(1))
        if ("下午" in text or "晚上" in text) and hour < 12:
            hour += 12
        if "凌晨" in text and hour == 12:
            hour = 0
        return max(0, min(23, hour)), 30
    dot_match = re.search(r"(\d{1,2})\s*点(?:\s*(\d{1,2}))?\s*分?", text)
    if dot_match:
        hour = int(dot_match.group(1))
        minute = int(dot_match.group(2)) if dot_match.group(2) else 0
        if ("下午" in text or "晚上" in text) and hour < 12:
            hour += 12
        if "凌晨" in text and hour == 12:
            hour = 0
        minute = max(0, min(59, minute))
        return max(0, min(23, hour)), minute
    return default_hour, default_minute


def _looks_like_schedule_message(message: str) -> bool:
    if not message:
        return False
    triggers = ["提醒", "定时", "闹钟", "日程", "定时任务", "通知", "叫我", "提示我"]
    time_hints = [
        "每天", "每周", "每月", "每年", "每小时", "每分钟", "明早", "今晚", "下周", "下个月",
        "早上", "上午", "中午", "下午", "晚上", "凌晨",
        "周一", "周二", "周三", "周四", "周五", "周六", "周日",
        "星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日",
        "工作日", "周末"
    ]
    task_verbs = ["写", "生成", "总结", "整理", "保存", "备份", "发送", "推送", "提醒", "通知", "检查", "同步", "下载", "导出", "还", "缴", "交", "支付", "复盘", "跟进"]
    has_time_hint = bool(
        any(item in message for item in time_hints)
        or re.search(r"\d{1,2}\s*点|\d{1,2}\s*[:：]\s*\d{2}", message)
        or re.search(r"\d+\s*(分钟|小时|天|周|个月|年)\s*(后|之后|以后)", message)
        or re.search(r"(半小时|一小时|一天|一周|一个月)\s*(后|之后|以后)", message)
    )
    has_date_hint = bool(
        re.search(r"\d{1,2}\s*月\s*\d{1,2}\s*(日|号)?", message)
        or re.search(r"(?<!\d)\d{1,2}\s*[\/\-.]\s*\d{1,2}(?!\d)", message)
    )
    if not has_time_hint and not has_date_hint:
        return False
    if _looks_like_birthday_reminder_message(message):
        return True
    # 含"保存到X盘"是明确的文件任务，有时间就是定时任务
    if re.search(r"保存到\s*[a-zA-Z]\s*盘", message):
        return True
    has_trigger = any(item in message for item in triggers)
    has_repeat_hint = any(item in message for item in time_hints)
    has_task_verb = any(item in message for item in task_verbs)
    return has_trigger or has_repeat_hint or has_task_verb


def _extract_reminder_content(message: str) -> str:
    text = message.strip()
    match = re.search(r"(提醒|记得|别忘了|闹钟|定时)\s*(我|我在|我去|我把|一下|一下子)?(.*)", text)
    if match and match.group(3).strip():
        content = match.group(3).strip()
    else:
        content = text
    content = re.sub(
        r"(每天|每周|每月|每年|每小时|每分钟|早上|上午|中午|下午|晚上|凌晨|今天|明天|后天|周一|周二|周三|周四|周五|周六|周日|星期一|星期二|星期三|星期四|星期五|星期六|星期日|工作日|周末)",
        " ",
        content
    )
    # 清理相对时间（10分钟后、2小时后、3天后、半小时后 等）
    content = re.sub(r"\d+\s*(分钟|小时|天|周|个月|年)\s*(后|之后|以后)", " ", content)
    content = re.sub(r"(半小时|一小时|一天|一周|一个月)\s*(后|之后|以后)", " ", content)
    content = re.sub(r"\d{1,2}\s*点半|\d{1,2}\s*点\s*\d{0,2}\s*分?|\d{1,2}\s*[:：]\s*\d{2}", " ", content)
    content = re.sub(r"[，。,.!！?？;；:：]", " ", content)
    content = " ".join(content.split())
    return content or text


def _extract_task_content(message: str) -> str:
    """从用户消息中提取纯任务内容，去掉时间相关的描述部分。

    例如：
      "5 分钟后，写首春天的诗，保存到 D:\\temp" → "写首春天的诗，保存到 D:\\temp"
      "每天早上 9 点总结今日新闻"              → "总结今日新闻"
      "明天上午10点写周报"                    → "写周报"
    """
    text = message.strip()

    # 0. 去掉"创建一个一次性定时任务："这类调度创建元指令，只保留未来要执行的正文。
    text = re.sub(
        r"^\s*(?:请|帮我|给我|麻烦你)?\s*"
        r"(?:创建|新建|添加|增加|设置|设定|建立|安排|生成)\s*"
        r"(?:一个|一条|1个)?\s*(?:一次性|单次|一次)?\s*"
        r"(?:定时任务|计划任务|任务|日程|提醒)\s*[:：,，、\s]*",
        "",
        text
    )

    # 1. 去掉句首的相对时间前缀（如"5分钟后，"、"半小时后 "）
    text = re.sub(
        r"^((\d+\s*(分钟|小时|天|周|个月|年)|半小时|一小时|一天|一周|一个月)\s*(后|之后|以后)[，,\s]*)+",
        "",
        text
    )

    # 2. 反复去除句首的时间词（每次去一个单元，循环直到不再变化）
    # 注意：更具体/更长的模式排在前面，避免被短模式抢先匹配
    _TIME_PREFIX_PATTERN = (
        r"^(?:"
        r"每月\d{1,2}[号日]?|"          # 每月X号
        r"每周[一二三四五六七日]|"        # 每周一/每周二...
        r"每[天周月年小时分钟]|"          # 每天/每周/每月等
        r"今天|明天|后天|工作日|周末|"
        r"周[一二三四五六七日]|星期[一二三四五六七日]|"
        r"早上|上午|中午|下午|晚上|凌晨|"
        r"\d{1,2}\s*[:：]\s*\d{2}|"     # HH:MM
        r"\d{1,2}\s*点半|"
        r"\d{1,2}\s*点\s*\d{1,2}\s*分|" # X点Y分
        r"\d{1,2}\s*(?:点钟?)|"          # X点/X点钟
        r"\d{1,2}\s*[号日](?!\d)"        # X号/X日
        r")\s*[，,\s]*"
    )
    prev = None
    while prev != text:
        prev = text
        text = re.sub(_TIME_PREFIX_PATTERN, "", text, count=1)
        text = text.lstrip("，,、 ")

    # 3. 去掉内嵌的相对时间表达（如任务中夹着的"X分钟后"）
    text = re.sub(r"\d+\s*(分钟|小时|天)\s*(后|之后|以后)", "", text)
    text = re.sub(r"(半小时|一小时|一天)\s*(后|之后|以后)", "", text)

    # 4. 清理多余空白和前导标点
    text = re.sub(r"^[，,、\s]+", "", text)
    text = " ".join(text.split())

    return text if text else message.strip()


def _find_duplicate_task(name: str, task_prompt: str, executor: str = "opencode") -> Optional[str]:
    """检查是否已存在相同或高度相似的定时任务。
    返回已有任务名称（如重复）或 None（不重复）。

    判定逻辑：
    - 任务名称完全相同 → 重复
    - task_prompt 去掉 __RUN_ONCE__/__REMINDER__ 标记后完全相同 → 重复
    - 名称相似度 > 0.8（SequenceMatcher）→ 重复
    """
    if not scheduler_router.scheduler:
        return None

    from difflib import SequenceMatcher

    def _strip_markers(text: str) -> str:
        """去掉 __RUN_ONCE__ / __REMINDER__ 前缀标记以比较实际内容"""
        for marker in ("__RUN_ONCE__\n", "__REMINDER__\n", "__RUN_ONCE__", "__REMINDER__"):
            text = text.replace(marker, "")
        return text.strip()

    existing_tasks = scheduler_router.scheduler.list_tasks()
    clean_prompt = _strip_markers(task_prompt)
    clean_name = name.strip()
    normalized_executor = _task_executor_from_target(executor)

    for t in existing_tasks:
        if _task_executor_from_target(getattr(t, "executor", "opencode")) != normalized_executor:
            continue
        # 1. 名称完全匹配
        if t.name.strip() == clean_name:
            return t.name
        # 2. prompt 内容完全匹配
        if _strip_markers(t.task_prompt) == clean_prompt and clean_prompt:
            return t.name
        # 3. 名称相似度高（> 0.8）
        ratio = SequenceMatcher(None, t.name.strip(), clean_name).ratio()
        if ratio > 0.8:
            return t.name

    return None


async def _try_create_scheduled_task(
    message: str,
    executor: str = "opencode",
    execution_model: str = "",
) -> Optional[str]:
    if not _looks_like_schedule_message(message):
        return None
    if not scheduler_router.scheduler:
        return None
    executor = _task_executor_from_target(executor)
    execution_model = _task_execution_model_from_chat_model(execution_model)
    # 从全局通知配置读取默认渠道
    _nc = app_config.notification
    _default_channels: list = []
    if _nc.app_enabled:
        _default_channels.append("app")
    if _nc.desktop_enabled:
        _default_channels.append("desktop")
    if _nc.lark_enabled:
        _default_channels.append("lark")
    if _nc.email_enabled:
        _default_channels.append("email")
    if not _default_channels:
        _default_channels = ["app"]
    try:
        if _looks_like_birthday_reminder_message(message):
            birthday = _extract_birthday_value(message or "")
            if birthday:
                md = re.search(r"(\d{1,2})月(\d{1,2})日", birthday)
                if md:
                    month = int(md.group(1))
                    day = int(md.group(2))
                    hour, minute = _extract_time_for_reminder(message)
                    subject = _extract_birthday_subject(message) or "我"
                    if subject == "我":
                        remind_text = "今天是你的生日，生日快乐！"
                        name = f"生日提醒：每年{month}月{day}日"
                    else:
                        remind_text = f"今天是{subject}的生日，记得送上祝福。"
                        name = f"生日提醒：{subject}每年{month}月{day}日"
                    task_prompt = f"__REMINDER__\n{remind_text}"
                    cron_expression = f"{minute} {hour} {day} {month} *"
                    # ---- 去重检查 ----
                    dup_name = _find_duplicate_task(name, task_prompt, executor=executor)
                    if dup_name:
                        logger.info(f"跳过重复生日提醒任务：'{name}'（已有：'{dup_name}'）")
                        return None
                    task = scheduler_router.scheduler.create_task(
                        name=name,
                        cron_expression=cron_expression,
                        task_prompt=task_prompt,
                        notify_channels=_default_channels,
                        executor=executor,
                        execution_model=execution_model,
                    )
                    next_run = task.next_run.isoformat() if task.next_run else "待计算"
                    return (
                        f"已创建定时任务：{task.name}\n"
                        f"Cron：{task.cron_expression}\n"
                        f"下次运行：{next_run}\n"
                        f"可在「定时任务」查看和管理。"
                    )
        cron_data = await scheduler_router.generate_cron_from_text(message)
        cron_expression = cron_data.get("cron") if cron_data else None
        if not cron_expression:
            return None
        run_once = (
            any(key in (message or "") for key in ["今天", "明天", "后天", "一次", "只提醒一次", "仅提醒一次", "提醒一次"])
            or bool(re.search(r"\d+\s*(分钟|小时|天)\s*(后|之后|以后)", message or ""))
            or bool(re.search(r"(半小时|一小时|一天)\s*(后|之后|以后)", message or ""))
        )
        is_reminder = any(key in (message or "") for key in ["提醒", "闹钟", "提示我", "叫我", "通知我", "别忘了", "记得"])
        if is_reminder:
            content = _extract_reminder_content(message)
            name = f"{'一次性' if run_once else ''}提醒：{content}" if content else ("一次性定时提醒" if run_once else "定时提醒")
            task_payload = f"提醒：{content}" if content else message.strip()
            task_prompt = f"__REMINDER__\n{task_payload}"
        else:
            # 提取纯任务内容（去掉时间相关描述，只保留真正要执行的任务）
            task_content = _extract_task_content(message)
            # 用于任务名称展示的简短版本（进一步去除标点和空白）
            action_text = re.sub(r"[，。,.!！?？;；:：]", " ", task_content)
            action_text = " ".join(action_text.split())
            if not action_text:
                action_text = message.strip()
            name = f"{'一次性' if run_once else ''}任务：{action_text}"
            # task_prompt 使用提取后的纯任务内容，而非原始消息
            # 这样执行时不会将"5分钟后"等时间描述混入任务指令
            task_prompt = task_content
            drive_match = re.search(r"保存到\s*([a-zA-Z])\s*盘", message)
            if drive_match:
                drive = drive_match.group(1).upper()
                output_dir = f"{drive}:\\codebot_tasks"
                task_prompt = (
                    f"{task_content}\n"
                    f"请将产出保存为 Markdown 文件到 {output_dir} 目录（如不存在请创建），"
                    f"文件名包含日期时间（例如 20260301_0800.md），并在完成后输出保存路径。"
                )

        if len(name) > 30:
            name = f"{name[:30]}..."
        if run_once:
            task_prompt = f"__RUN_ONCE__\n{task_prompt}"
        # ---- 去重检查 ----
        dup_name = _find_duplicate_task(name, task_prompt, executor=executor)
        if dup_name:
            logger.info(f"跳过重复定时任务：'{name}'（已有：'{dup_name}'）")
            return None
        task = scheduler_router.scheduler.create_task(
            name=name,
            cron_expression=cron_expression,
            task_prompt=task_prompt,
            notify_channels=_default_channels,
            executor=executor,
            execution_model=execution_model,
        )
        next_run = task.next_run.isoformat() if task.next_run else "待计算"
        return f"已创建定时任务：{task.name}\nCron：{task.cron_expression}\n下次运行：{next_run}\n可在“定时任务”查看和管理。"
    except Exception as e:
        logger.error(f"创建定时任务失败: {e}")
        return None


def _get_chat_memory_manager() -> MemoryManager:
    global chat_memory_manager
    if chat_memory_manager is None:
        chat_memory_manager = MemoryManager()
    return chat_memory_manager


async def _save_memory_content(content: str, category: str, raw_message: str) -> bool:
    selected_skill, cleaned_message, opencode_skill_fallback = _extract_requested_skill(raw_message)
    message_for_context = cleaned_message or raw_message
    manager = _get_chat_memory_manager()
    metadata = {"source": "chat", "raw": raw_message}
    try:
        await manager.save_long_term_memory(content=content, category=category, metadata=metadata)
        return True
    except Exception as e:
        try:
            cursor = manager.sqlite_db.cursor()
            cursor.execute(
                "SELECT id FROM long_term_memories WHERE category = ? AND content = ? ORDER BY id DESC LIMIT 1",
                (category, content)
            )
            exists = cursor.fetchone()
            if not exists:
                cursor.execute(
                    "INSERT INTO long_term_memories (category, content, metadata) VALUES (?, ?, ?)",
                    (category, content, json.dumps(metadata, ensure_ascii=False))
                )
                manager.sqlite_db.commit()
            return True
        except Exception as inner:
            logger.error(f"保存记忆失败: {inner}")
            return False
        finally:
            logger.error(f"记忆向量索引失败: {e}")


def _looks_like_memory_message(message: str) -> bool:
    if not message:
        return False
    text = message.strip()
    # 如果消息明显是定时任务（含时间+任务动词），不识别为记忆
    if _looks_like_schedule_message(text):
        return False
    triggers = [
        "帮我记住", "请记住", "记住", "记一下", "记下",
        "写入记忆", "保存记忆", "保存到记忆", "存为记忆", "存进记忆",
        "更新记忆", "修改记忆", "更正记忆"
    ]
    # "保存到X盘" 是文件保存，不是记忆
    if re.search(r"保存到\s*[a-zA-Z]\s*盘", text):
        return False
    return any(trigger in text for trigger in triggers)


def _looks_like_birthday_memory_intent(message: str) -> bool:
    text = (message or "").strip()
    if not text:
        return False
    if "生日" not in text:
        return False
    value = _extract_birthday_value(text)
    if not value:
        return False
    if any(key in text for key in ["记住", "保存", "写入记忆", "保存记忆", "更新记忆", "修改记忆", "更正记忆"]):
        return True
    if _looks_like_memory_message(text):
        return True
    if re.search(r"^(我的生日|我生日)\s*(是|为|改为|改成)", text):
        return True
    if any(key in text for key in ["更新", "修改", "更正", "改为", "改成"]):
        return True
    if any(key in text for key in ["吗", "？", "?", "几号", "哪天", "什么时候"]):
        return False
    if re.search(r"(我的|我|姐姐|哥哥|妈妈|爸爸|父母|老婆|妻子|老公|女儿|儿子|朋友|同事|领导|老师|老板).{0,8}(生日)", text) and ("是" in text or "为" in text):
        return True
    match = re.search(r"(.{1,12})的生日\s*(是|为)", text)
    if match:
        subject = (match.group(1) or "").strip()
        if subject and subject not in ["什么", "谁", "哪天", "几号"]:
            return True
    return False


def _extract_birthday_subject(message: str) -> Optional[str]:
    text = (message or "").strip()
    if not text:
        return None
    if re.search(r"(我的生日|我生日)", text):
        return "我"
    match = re.search(r"是\s*([^，。,.!！?？;；:：\s]{1,12})\s*的生日", text)
    if match:
        subject = (match.group(1) or "").strip()
        if subject and (not re.search(r"[\d月日号]", subject)):
            return subject
    match = re.search(r"([^，。,.!！?？;；:：\s]{1,12})的生日", text)
    if match:
        subject = (match.group(1) or "").strip()
        if subject and (not re.search(r"[\d月日号]", subject)):
            return subject
    for candidate in ["姐姐", "哥哥", "妈妈", "爸爸", "父母", "老婆", "妻子", "老公", "女儿", "儿子"]:
        if candidate in text and "生日" in text:
            return candidate
    return None


def _extract_birthday_value(message: str) -> Optional[str]:
    text = (message or "").strip()
    if not text:
        return None

    match = re.search(r"(\d{1,2})\s*月\s*(\d{1,2})\s*(日|号)?", text)
    if match:
        month = int(match.group(1))
        day = int(match.group(2))
        if 1 <= month <= 12 and 1 <= day <= 31:
            return f"{month}月{day}日"

    match = re.search(r"(?<!\d)(\d{1,2})\s*[\/\-.]\s*(\d{1,2})(?!\d)", text)
    if match:
        month = int(match.group(1))
        day = int(match.group(2))
        if 1 <= month <= 12 and 1 <= day <= 31:
            return f"{month}月{day}日"

    return None


def _is_birthday_question(message: str) -> bool:
    text = (message or "").strip()
    if not text:
        return False
    if "生日" in text:
        return True
    if any(key in text for key in ["出生日期", "生日是哪天", "生日是什么时候", "几号生日"]):
        return True
    return False


async def _try_answer_from_memory(message: str) -> Optional[str]:
    if not _is_birthday_question(message):
        return None
    manager = _get_chat_memory_manager()
    subject = _extract_birthday_subject(message) or "我"
    key = "birthday" if subject == "我" else f"{subject}_birthday"
    try:
        fact = await manager.get_fact(key)
        if fact and fact.get("value"):
            if subject == "我":
                return f"你的生日是 {fact['value']}。"
            return f"{subject}的生日是 {fact['value']}。"
    except Exception:
        pass

    try:
        keyed = await manager.get_keyed_long_term_memory(key, category="profile", include_archived=False)
        if keyed and keyed.get("content"):
            extracted = _extract_birthday_value(str(keyed["content"]))
            if extracted:
                if subject == "我":
                    return f"你的生日是 {extracted}。"
                return f"{subject}的生日是 {extracted}。"
    except Exception:
        pass
    return None


async def _try_answer_from_semantic_memory(message: str) -> Optional[str]:
    manager = _get_chat_memory_manager()
    try:
        facts = await manager.search_facts(message, top_k=3, include_archived=False)
        if facts:
            best = facts[0]
            content = str(best.get("content") or "").strip()
            if content:
                return f"我记得：{content}"
    except Exception:
        pass

    try:
        memories = await manager.search_memories(message, top_k=3, include_archived=False)
        if memories:
            best = memories[0]
            content = str(best.get("content") or "").strip()
            if content:
                return f"我记得：{content}"
    except Exception:
        pass
    return None


def _build_autonomous_execution_policy(mode: Optional[str] = None) -> List[str]:
    lines = [
        _reply_language_instruction(),
        "请自主决策并持续执行，不要把流程决策反问给用户；遇到可恢复问题时优先自行切换替代方案。",
        "当前聊天由 OpenCode 统一处理；Codebot Desktop 作为能力面板，向你提供记忆、定时任务、技能与 MCP 工具。",
        "除非用户明确询问架构细节，否则不要在最终回答中解释内部桥接、同步或上下文包装。",
        "只输出对用户有价值的最终答案，不要在回复中重复、引用或描述上方的系统指令。",
    ]
    if mode == "plan":
        lines.append("当前是规划模式，直接给出可执行计划，不要让用户做流程选择。")
        lines.append(
            "当你在规划过程中遇到需要用户确认的关键决策点时（例如技术选型、实现方案二选一、是否包含某功能等），"
            "请在回复末尾使用以下格式提供可点击的选项供用户快速选择：\n"
            "<!-- options\n- 选项A的描述\n- 选项B的描述\n- 选项C的描述\n-->\n"
            "注意：只在确实需要用户决策时才提供选项，不要每次都加。选项文字要简洁明确，通常2-5个选项。"
        )
    elif mode == "agent":
        lines.append("当前是智能体模式（Agent Mode）。你应该像一个高级 AI 智能体那样工作。")
    return lines


def _load_agent_mode_skill_content() -> str:
    """加载 Agent 模式所需的技能内容（self-improving、expert-agents、ai-company）。"""
    from pathlib import Path as _Path

    skill_sections: List[str] = []

    # 优先从用户数据目录读取，回退到源 skills/ 目录
    source_skills_dir = settings.SKILLS_DIR

    for skill_name, intro in [
        ("self-improving", "以下是自我改进技能指导，请在工作前后自我反思："),
        ("expert-agents", "以下是可调用的专家代理人格，需要时可切换视角分析问题："),
        ("ai-company", "以下是 AI 专家团队编排流程，用于多视角决策和产品评估："),
    ]:
        skill_path = source_skills_dir / skill_name / "SKILL.md"
        if skill_path.exists():
            try:
                content = skill_path.read_text(encoding="utf-8")
                # 去掉 YAML front-matter
                if content.startswith("---"):
                    end_idx = content.find("---", 3)
                    if end_idx != -1:
                        content = content[end_idx + 3:].strip()
                # 截断过长内容（每个技能最多2000字符）
                if len(content) > 2000:
                    content = content[:2000] + "\n... (内容已截断)"
                skill_sections.append(f"## {skill_name}\n{intro}\n{content}")
            except Exception:
                pass

    return "\n\n".join(skill_sections)


def _extract_requested_skill(message: str) -> Tuple[Optional[dict], str, bool]:
    """Resolve an explicitly requested skill using Codebot's priority order."""
    text = message or ""
    registry = get_skill_registry()

    marker = re.search(r"使用技能\[([^\]]+)\]", text)
    if marker:
        skill_id = marker.group(1).strip()
        item = registry.find(skill_id)
        cleaned = (text[:marker.start()] + text[marker.end():]).strip()
        return item, cleaned or text, bool(item and item.get("source") == OPENCODE)

    slash = re.search(r"(?im)^\s*/skill\s+([^\n]+)", text)
    if slash:
        query = slash.group(1).strip()
        item = registry.find_by_query(query, allow_opencode=True)
        cleaned = (text[:slash.start()] + text[slash.end():]).strip()
        return item, cleaned or text, bool(item and item.get("source") == OPENCODE)

    natural = re.search(r"(?:使用|调用|启用)\s+(.{1,80}?)\s+技能", text)
    if natural:
        query = natural.group(1).strip()
        item = registry.find_by_query(query, allow_opencode=True)
        return item, text, bool(item and item.get("source") == OPENCODE)

    return None, text, False


def _build_skill_system_context(skill: dict) -> str:
    source = skill.get("source") or ""
    source_label = skill.get("sourceLabel") or skill.get("source_label") or source
    content = skill.get("skill_md_content") or ""
    if not content and skill.get("skill_md_path"):
        try:
            content = Path(skill.get("skill_md_path")).read_text(encoding="utf-8")
        except Exception:
            content = ""
    if len(content) > 12000:
        content = content[:12000] + "\n\n..."
    return (
        f"=== Codebot skill selected: {skill.get('name') or skill.get('slug')} ({source_label}) ===\n"
        "Use the following SKILL.md as the primary workflow for this turn. "
        "If it conflicts with the user's latest explicit request, follow the user request.\n\n"
        f"{content}"
    )


async def _build_opencode_prompt_parts(
    message: str,
    mode: Optional[str] = None,
    project_dir: Optional[str] = None,
    target: Optional[str] = None,
    knowledge_paths: Optional[List[str]] = None,
    model: Optional[str] = None,
) -> Tuple[str, str]:
    """
    构建 OpenCode 提示词，返回 (system_prompt, user_message) 元组。

    system_prompt 通过 OpenCode API 的独立 system 字段传递，不混入用户消息文本，
    从根本上解决 AI 把系统指令原文回显到回复中的问题。
    """
    raw_message = message or ""
    try:
        selected_skill, cleaned_message, opencode_skill_fallback = _extract_requested_skill(raw_message)
    except Exception as exc:
        logger.debug(f"[skill] 解析请求技能失败: {exc}")
        selected_skill, cleaned_message, opencode_skill_fallback = None, raw_message, False
    # 用于记忆/事实检索的查询文本（去除技能标记）
    message_for_context = cleaned_message or raw_message

    manager = _get_chat_memory_manager()
    facts_context: List[str] = []
    memories_context: List[str] = []
    habit_context: List[str] = []
    preference_context: List[str] = []
    profile_context: List[str] = []

    # ── 事实记忆（结构化 key-value）──────────────────────────────────────────
    try:
        facts = await manager.search_facts(message_for_context, top_k=5, include_archived=False)
        for item in facts:
            content = str(item.get("content") or "").strip()
            if content and content not in facts_context:
                facts_context.append(content)
    except Exception:
        pass

    # ── 向量语义记忆（跨所有分类）────────────────────────────────────────────
    try:
        memories = await manager.search_memories(message_for_context, top_k=5, include_archived=False)
        for item in memories:
            content = str(item.get("content") or "").strip()
            cat = str(item.get("category") or "")
            if not content:
                continue
            if cat == "habit" and content not in habit_context:
                habit_context.append(content)
            elif cat == "preference" and content not in preference_context:
                preference_context.append(content)
            elif cat == "profile" and content not in profile_context:
                profile_context.append(content)
            elif content not in memories_context:
                memories_context.append(content)
    except Exception:
        pass

    # ── 专项分类检索（补充语义检索未命中的内容）──────────────────────────────
    for cat, target_bucket in [
        ("habit", habit_context),
        ("preference", preference_context),
        ("profile", profile_context),
    ]:
        if len(target_bucket) < 3:
            try:
                extra = await manager.search_memories(
                    message_for_context, top_k=3, category=cat, include_archived=False
                )
                for item in extra:
                    content = str(item.get("content") or "").strip()
                    if content and content not in target_bucket:
                        target_bucket.append(content)
            except Exception:
                pass

    policy_lines: List[str] = _build_autonomous_execution_policy(mode=mode)
    memory_lines: List[str] = []
    has_any = any([facts_context, habit_context, preference_context, profile_context, memories_context])
    if facts_context:
        memory_lines.append("【用户事实记忆（可信，优先使用；如有冲突以最新为准）】")
        for item in facts_context[:5]:
            memory_lines.append(f"- {item}")
    if profile_context:
        memory_lines.append("【用户个人信息】")
        for item in profile_context[:3]:
            memory_lines.append(f"- {item}")
    if preference_context:
        memory_lines.append("【用户偏好】")
        for item in preference_context[:3]:
            memory_lines.append(f"- {item}")
    if habit_context:
        memory_lines.append("【用户习惯】")
        for item in habit_context[:3]:
            memory_lines.append(f"- {item}")
    if memories_context:
        memory_lines.append("【用户长期记忆（可信，尽量参考）】")
        for item in memories_context[:5]:
            memory_lines.append(f"- {item}")

    # ── 构建 system prompt（仅含系统指令，不含用户消息）─────────────────────
    system_lines: List[str] = ["你正在 OpenCode 中处理用户消息。"]
    system_lines.extend(policy_lines)
    if _looks_like_codebot_schedule_creation_request(raw_message):
        system_lines.append(_build_codebot_scheduler_boundary(target))
    if _looks_like_skill_creation_intent(raw_message):
        system_lines.append(
            "本轮消息来自 Codebot 聊天，且用户有创建、生成、保存或沉淀 skill 的意图。"
            "如果需要创建技能，请仍按 OpenCode 原有 skill-creator/agent-skill 逻辑完成，"
            "但不要把 skill 保存到用户的文件存储路径；Codebot 会在任务结束后把本轮新生成的 skill 迁移为 Codebot 的自动生成技能。"
        )
    if has_any:
        system_lines.append("以下是与当前问题相关的用户记忆，请在回答中参考；若与用户本轮消息冲突，以用户本轮消息为准。")
        system_lines.extend(memory_lines)

    if selected_skill and selected_skill.get("source") in {AUTO_GENERATED, BUILTIN, EXTERNAL, HERMES, OPENCLAW}:
        system_lines.append(_build_skill_system_context(selected_skill))

    normalized_target = (target or "codebot").strip().lower()
    if normalized_target == "hermes":
        system_lines.append(_build_hermes_context(model=model))
    elif normalized_target == "obsidian":
        obsidian_context = _build_obsidian_context(cleaned_message or raw_message, knowledge_paths)
        if obsidian_context:
            system_lines.append(obsidian_context)
        else:
            system_lines.append(
                "Obsidian mode is active, but no configured Obsidian vault or knowledge base was selected. "
                "Ask the user to configure Obsidian Settings if Markdown vault access is required."
            )

    # ── 注入文件存储/项目目录（项目目录优先）──────────────────────────────
    if project_dir and project_dir.strip():
        system_lines.append(
            f"用户当前工作项目目录为：{project_dir.strip()}。"
            f"所有文件操作（读取、搜索、创建、编辑、保存、导出）请基于此项目目录进行，"
            f"除非用户明确指定其他路径。"
        )
    else:
        try:
            from config import app_config as _cfg
            file_storage_path = (_cfg.general.file_storage_path or "").strip()
            if file_storage_path:
                system_lines.append(
                    f"用户指定的文件存储目录为：{file_storage_path}。"
                    f"当你需要生成、保存或导出文件时（如 MD、CSV、TXT、HTML 等），"
                    f"请将文件保存到此目录下，无需再询问用户保存位置。"
                )
        except Exception:
            pass

    # ── Agent 模式：注入自我改进、专家代理、AI 团队技能 ───────────────────────
    if mode == "agent":
        agent_skills = _load_agent_mode_skill_content()
        if agent_skills:
            system_lines.append(
                "=== Agent 模式技能指导 ===\n"
                "你现在处于智能体（Agent）模式。请遵循以下技能指导来提升工作质量：\n"
                "1. 开始工作前先自我反思：回顾相关记忆、评估任务复杂度\n"
                "2. 完成工作后自我批评：检查潜在问题、评估输出质量\n"
                "3. 需要多视角分析时，调用专家代理获取不同领域的意见\n"
                "4. 将学到的经验教训记录到记忆中，持续改进\n\n"
                f"{agent_skills}"
            )

    system_prompt = "\n\n".join(system_lines)
    if selected_skill and opencode_skill_fallback:
        cleaned_message = f"/skill {selected_skill.get('slug')}\n\n{cleaned_message or raw_message}"
    # user_message 只含纯净的用户输入，不混入任何系统指令
    return system_prompt, cleaned_message or raw_message


async def _build_opencode_prompt_with_memory(message: str, mode: Optional[str] = None) -> str:
    """兼容旧调用：返回拼合后的单字符串（system + user）。内部已不再使用，保留供外部兼容。"""
    system_prompt, user_message = await _build_opencode_prompt_parts(message, mode=mode)
    return f"{system_prompt}\n\n【用户消息】{user_message}"

def _extract_memory_content(message: str) -> str:
    text = (message or "").strip()
    if not text:
        return ""

    quoted_match = re.search(r"[“\"'‘](.+?)[”\"'’]", text)
    if quoted_match and quoted_match.group(1).strip():
        candidate = quoted_match.group(1).strip()
    else:
        match = re.search(
            r"(帮我|请|麻烦)?(记住|保存|记一下|记下|存一下|存下)\s*(一下|下)?\s*[:：]?\s*(.*)",
            text
        )
        candidate = (match.group(4) if match else text).strip()

    candidate = re.sub(r"(这个|该)?(地址|位置|地点|信息|内容)\s*$", "", candidate).strip()
    candidate = re.sub(r"[，。,.!！?？;；]\s*$", "", candidate).strip()
    return candidate


def _guess_memory_category(message: str, content: str) -> str:
    text = f"{message} {content}".strip()
    # ── preference（偏好）── 优先匹配，避免被 profile 抢走
    if any(key in text for key in [
        "喜欢", "偏好", "偏爱", "偏向", "倾向", "首选", "爱用", "不喜欢",
        "讨厌", "不想", "不爱", "不习惯", "风格", "方式", "模式", "格式",
        "回复风格", "编程语言", "框架", "工具", "编辑器", "IDE", "主题",
    ]):
        return "preference"
    # ── habit（习惯）──
    if any(key in text for key in [
        "习惯", "通常", "一般", "平时", "经常", "常常", "总是", "每天",
        "每次", "常用", "惯用", "日常", "作息", "规律",
    ]):
        return "habit"
    # ── contact（联系方式）── 必须在 profile 之前，关键词不会误伤
    if any(key in text for key in ["电话", "手机号", "联系方式", "号码", "微信", "邮箱"]):
        return "contact"
    # ── address（地址）──
    if any(key in text for key in ["地址", "住址", "位置", "地点"]):
        return "address"
    # ── profile（个人信息）── 仅限真正的身份/人口学信息
    if any(key in text for key in [
        "生日", "个人信息", "姓名", "名字", "年龄", "身份", "身份证",
        "账号", "账户", "密码", "口令", "职业", "工作", "学校", "职位",
    ]):
        return "profile"
    return "note"


async def _try_save_memory(message: str) -> Optional[str]:
    manager = _get_chat_memory_manager()
    birthday_value = _extract_birthday_value(message)
    if birthday_value and _looks_like_birthday_memory_intent(message):
        subject = _extract_birthday_subject(message) or "我"
        key = "birthday" if subject == "我" else f"{subject}_birthday"
        metadata = {"source": "chat", "raw": message, "memory_key": key, "subject": subject}
        try:
            await manager.upsert_fact(key=key, value=birthday_value, metadata=metadata)
        except Exception:
            pass
        try:
            await manager.upsert_keyed_long_term_memory(
                memory_key=key,
                content=f"{'生日' if subject == '我' else subject + '的生日'}：{birthday_value}",
                category="profile",
                metadata=metadata
            )
        except Exception:
            return "记忆保存失败，请稍后重试"
        if subject == "我":
            return f"我已记录你的生日是 {birthday_value}。"
        return f"我已记录{subject}的生日是 {birthday_value}。"

    if not _looks_like_memory_message(message):
        return None

    content = _extract_memory_content(message)
    if not content:
        return None

    category = _guess_memory_category(message, content)
    metadata = {"source": "chat", "raw": message}

    try:
        await manager.save_long_term_memory(content=content, category=category, metadata=metadata)
    except Exception as e:
        try:
            cursor = manager.sqlite_db.cursor()
            cursor.execute(
                "SELECT id FROM long_term_memories WHERE category = ? AND content = ? ORDER BY id DESC LIMIT 1",
                (category, content)
            )
            exists = cursor.fetchone()
            if not exists:
                cursor.execute(
                    "INSERT INTO long_term_memories (category, content, metadata) VALUES (?, ?, ?)",
                    (category, content, json.dumps(metadata, ensure_ascii=False))
                )
                manager.sqlite_db.commit()
        except Exception as inner:
            logger.error(f"保存记忆失败: {inner}")
            return "记忆保存失败，请稍后重试"
        logger.error(f"记忆向量索引失败: {e}")

    return f"好的，我已经记住了：{content}\n可在“记忆”页面查看。"


def _build_skill_from_conversation(messages: List[dict], request: SkillGenerateRequest) -> dict:
    user_messages = [item.get("content") for item in messages if item.get("role") == "user"]
    last_user_message = ""
    for item in reversed(user_messages):
        if item and str(item).strip():
            last_user_message = str(item).strip()
            break
    fallback_name = generate_conversation_title(last_user_message or "对话技能")
    description = (request.description or last_user_message).strip()
    if len(description) > 200:
        description = f"{description[:200]}..."
    return {
        "name": (request.name or fallback_name or "未命名技能").strip(),
        "description": description,
        "version": request.version or "2.0.0",
        "source": request.source or "chat",
        "enabled": bool(request.enabled)
    }


def _should_materialize_skill(user_message: str, assistant_response: str, conversation_id: Optional[int] = None) -> bool:
    user_text = (user_message or "").strip()
    answer_text = _sanitize_assistant_output(assistant_response or "")
    if _skill_content_is_noise(answer_text):
        return False
    if not user_text or len(user_text) < 8:
        return False
    if re.fullmatch(r"(你好|您好|在吗|hi|hello|hey)[!！。,. ]*", user_text, flags=re.IGNORECASE):
        return False
    direct_skill_request = any(word in user_text for word in ["生成skill", "生成技能", "沉淀为技能", "保存为技能", "做成技能", "以后遇到", "下次遇到"])
    if len(answer_text) < (120 if direct_skill_request else 180):
        return False
    structure_count = len(re.findall(r"(?:\n\s*[-*]|\n\s*\d+[.)、])", answer_text))
    has_code_or_sections = bool(re.search(r"```|\n#{2,}\s+|\n(?:步骤|流程|命令|注意事项|验证)[:：]", answer_text))
    if not direct_skill_request and structure_count < 1 and not has_code_or_sections:
        return False
    trigger_words = [
        "步骤", "流程", "脚本", "命令", "自动化", "排查", "修复", "部署", "配置", "方案", "实现", "改造",
        "workflow", "pipeline", "script", "troubleshoot", "deploy", "automation", "refactor", "migration"
    ]
    hit = sum(1 for w in trigger_words if w.lower() in f"{user_text}\n{answer_text}".lower())
    if direct_skill_request:
        return True
    if hit < 2:
        return False
    if conversation_id is None:
        return hit >= 2 and len(answer_text) >= 260

    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        history = memory_manager.sqlite_db.cursor().execute(
            "SELECT content FROM messages WHERE conversation_id = ? AND role = 'user' ORDER BY created_at DESC LIMIT 12",
            (conversation_id,)
        ).fetchall()
    except Exception:
        history = []

    normalized_user = re.sub(r"\s+", " ", user_text.lower()).strip()
    repeated_hits = 0
    for row in history:
        prev = str(row[0] or "").strip()
        prev_norm = re.sub(r"\s+", " ", prev.lower()).strip()
        if not prev_norm or prev_norm == normalized_user:
            continue
        if _skill_name_similar(prev_norm, normalized_user):
            repeated_hits += 1
            if repeated_hits >= 1:
                return True

    return False


def _skill_name_similar(a: str, b: str) -> bool:
    """粗略相似度：两个字符串共享超过 60% 的字符集，或一方是另一方的子串。"""
    a, b = (a or "").strip().lower(), (b or "").strip().lower()
    if not a or not b:
        return False
    if a in b or b in a:
        return True
    set_a, set_b = set(a), set(b)
    if not set_a or not set_b:
        return False
    overlap = len(set_a & set_b) / min(len(set_a), len(set_b))
    return overlap >= 0.6


async def _materialize_reusable_skill(user_message: str, assistant_response: str, conversation_id: Optional[int] = None) -> bool:
    cleaned_response = _sanitize_assistant_output(assistant_response or "")
    if not _should_materialize_skill(user_message, cleaned_response, conversation_id=conversation_id):
        return False
    title = generate_conversation_title(user_message or "自动技能")
    item = await _materialize_skill_content(
        name=f"自动技能-{title}",
        description=cleaned_response.strip().replace("\n", " "),
        user_message=user_message,
        assistant_response=cleaned_response,
        slug_hint=title,
    )
    return bool(item)


async def _async_create_skill_via_opencode(session_id: str, prompt: str) -> None:
    """异步通过 OpenCode 的 skill-creator 技能创建新技能。"""
    try:
        if opencode_ws and getattr(opencode_ws, "connected", False):
            snapshot = capture_opencode_skill_snapshot()
            import time
            started_at = time.time()
            await opencode_ws.send_message(
                session_id=session_id,
                message=prompt,
            )
            logger.info(f"[skill] skill-creator 异步任务完成: {session_id}")
            migrate_new_opencode_skills_to_codebot(
                snapshot=snapshot,
                since=started_at,
                reason="codebot_async_skill_creator",
            )
    except Exception as exc:
        logger.warning(f"[skill] skill-creator 异步任务失败: {exc}")


async def _rescue_misplaced_skills() -> None:
    """Best-effort migration for recently created OpenCode skills."""
    import time

    migrate_new_opencode_skills_to_codebot(
        snapshot=None,
        since=time.time() - 60,
        reason="codebot_recent_skill_rescue",
    )


def _write_auto_skill_md(skill_md_path, name: str, description: str,
                          user_message: str, assistant_response: str) -> None:
    """将自动技能写入 SKILL.md 文件（builtin 格式，可编辑）。"""
    write_auto_skill_md(
        Path(skill_md_path),
        name=name,
        description=description,
        body=assistant_response,
        user_message=user_message,
        slug=Path(skill_md_path).parent.name,
    )
    return
    # 截断过长内容
    content_body = assistant_response.strip()
    if len(content_body) > 2000:
        content_body = content_body[:2000] + "\n\n..."
    skill_content = f"""---
name: {name}
description: {description[:120] if len(description) > 120 else description}
---

# {name}

## 技能概述

此技能由对话自动生成，记录了以下场景的处理方式：

> {user_message[:200] if len(user_message) > 200 else user_message}

## 主要内容

{content_body}
"""
    skill_md_path.write_text(skill_content, encoding="utf-8")


def _write_skill(skill_id: str, data: dict):
    settings.SKILLS_DIR.mkdir(parents=True, exist_ok=True)
    path = settings.SKILLS_DIR / f"{skill_id}.json"
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


@router.post("/conversations", response_model=MessageResponse)
async def create_conversation(
    title: str = Body("新对话"),
    project_dir: Optional[str] = Body(None),
    conversation_type: str = Body("normal")
):
    """创建对话"""
    try:
        # 初始化数据库连接
        conversations_db.connect()
        
        memory_manager = MemoryManager()
        if conversation_type == "multi_agent_hub":
            conversation = await memory_manager.ensure_multi_agent_hub()
            return MessageResponse(
                success=True,
                data=conversation,
                message="多Agent群聊已就绪"
            )

        conversation_id = await memory_manager.create_conversation(title, project_dir=project_dir)
        
        return MessageResponse(
            success=True,
            data={"id": conversation_id, "title": title, "project_dir": project_dir},
            message="对话创建成功"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/conversations")
async def list_conversations(
    limit: int = 50,
    offset: int = 0,
    archived: bool = False
):
    """获取对话列表"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        
        if not archived:
            await memory_manager.ensure_multi_agent_hub()

        conversations = await memory_manager.get_conversations(
            limit=limit,
            offset=offset,
            archived=archived
        )
        
        return {
            "success": True,
            "data": {
                "items": conversations,
                "total": len(conversations)
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/conversations/{conversation_id}")
async def get_conversation(conversation_id: int):
    """获取对话详情"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        
        conversation = await memory_manager.get_conversation(conversation_id)
        
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        
        return {
            "success": True,
            "data": conversation
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.patch("/conversations/{conversation_id}/project_dir")
async def update_conversation_project_dir(
    conversation_id: int,
    project_dir: Optional[str] = Body(None, embed=True)
):
    """更新对话关联的项目目录"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        await memory_manager.update_conversation_project_dir(conversation_id, project_dir)
        return {"success": True, "data": {"project_dir": project_dir}, "message": "项目目录已更新"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/conversations/{conversation_id}")
async def delete_conversation(conversation_id: int):
    """删除对话"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        if _is_multi_agent_hub(conversation):
            raise HTTPException(status_code=400, detail="多Agent群聊不能删除，请使用清空")

        await memory_manager.delete_conversation(conversation_id)
        
        return {
            "success": True,
            "message": "对话已删除"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.patch("/conversations/{conversation_id}/title")
async def update_conversation_title(conversation_id: int, request: UpdateTitleRequest):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        await memory_manager.update_conversation_title(conversation_id, request.title)
        return {
            "success": True,
            "message": "标题已更新"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/pin")
async def toggle_conversation_pinned(conversation_id: int, request: TogglePinnedRequest):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        await memory_manager.set_conversation_pinned(conversation_id, request.pinned)
        return {
            "success": True,
            "message": "置顶状态已更新"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/archive")
async def toggle_conversation_archived(conversation_id: int, request: ToggleArchiveRequest):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        await memory_manager.set_conversation_archived(conversation_id, request.archived)
        return {
            "success": True,
            "message": "归档状态已更新"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/group")
async def toggle_conversation_group(conversation_id: int, request: ToggleGroupRequest):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        if _is_multi_agent_hub(conversation):
            raise HTTPException(status_code=400, detail="多Agent群聊工作台不能退出群聊")
        role = request.group_role or conversation.get("group_role") or conversation.get("title")
        await memory_manager.set_conversation_group(conversation_id, request.is_group, role)
        return {
            "success": True,
            "message": "群聊状态已更新"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/clear")
async def clear_conversation(conversation_id: int, request: ClearConversationRequest):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        if not request.confirm:
            raise HTTPException(status_code=400, detail="需要确认清空")
        await memory_manager.clear_conversation_messages(conversation_id)
        return {"success": True, "message": "对话内容已清空"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/multi-agent/members")
async def list_multi_agent_members(project_dir: Optional[str] = None):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        await memory_manager.ensure_multi_agent_hub()
        members = await memory_manager.get_multi_agent_members(project_dir=project_dir)
        return {"success": True, "data": {"items": members, "total": len(members)}}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/multi-agent/{hub_id}/dispatch")
async def dispatch_multi_agent_task(hub_id: int, request: MultiAgentDispatchRequest):
    hub_key = str(hub_id)
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        hub = await memory_manager.get_conversation(hub_id)
        if not _is_multi_agent_hub(hub):
            raise HTTPException(status_code=400, detail="当前对话不是多Agent群聊")

        members = await memory_manager.get_multi_agent_members(project_dir=request.project_dir)
        if not members:
            content = "请先把至少一个普通对话加入多Agent群聊，并为它设置角色（如前端、后端、数据库、测试）。"
            await memory_manager.save_message(hub_id, "assistant", content)
            return {"success": True, "data": {"content": content, "results": []}}

        steps = _plan_multi_agent_tasks(request.message, members)
        plan_text = _format_multi_agent_plan(request.message, steps)
        progress_lines = [plan_text, "", "## 执行过程"]
        _runtime_start(hub_key)
        mark_conversation_running(hub_key)
        _multi_agent_dispatch_state[hub_key] = {
            "member_ids": [str(assignment["member"].get("id")) for step in steps for assignment in step],
            "aborted": False,
            "updated_at": datetime.now(),
        }
        _runtime_set_content(hub_key, "\n".join(progress_lines).strip())
        _runtime_append_event(hub_key, {"type": "status", "message": "已生成多Agent任务计划"})

        results: List[Dict[str, str]] = []
        upstream_outputs: List[Dict[str, str]] = []
        for step_index, step in enumerate(steps, 1):
            state = _multi_agent_dispatch_state.get(hub_key) or {}
            if state.get("aborted"):
                raise asyncio.CancelledError()
            _append_hub_progress(hub_key, progress_lines, f"### 第 {step_index} 步开始：{'并行' if len(step) > 1 else '串行'}")

            for assignment in step:
                member = assignment["member"]
                role = _conversation_role_label(member)
                task_text = assignment["task"]
                upstream_text = "\n\n".join(
                    f"【上游产物 - {item['role']}】\n{item['reply']}" for item in upstream_outputs
                )
                delegated_message = (
                    f"【多Agent群聊分配任务】\n"
                    f"你的角色：{role}\n"
                    f"总任务：{request.message}\n"
                    f"当前步骤：第 {step_index} 步\n"
                    f"你负责的子任务：{task_text}\n"
                )
                if assignment.get("depends_on_previous") and upstream_text:
                    delegated_message += f"\n上游 Agent 已完成的产物如下，请基于它继续处理：\n{upstream_text}\n"
                delegated_message += "\n请只处理与你角色相关的部分，完成后说明结果、修改点、风险和需要其他 Agent 配合的事项。"

                _append_hub_progress(hub_key, progress_lines, f"- 分配给 {role}（对话 #{member['id']}）：{task_text}")
                await memory_manager.save_message(member["id"], "user", delegated_message)
                reply = await _execute_opencode(
                    delegated_message,
                    model=request.model,
                    mode=request.mode or "agent",
                    conversation_id=str(member["id"]),
                    project_dir=member.get("project_dir") or request.project_dir,
                )
                state = _multi_agent_dispatch_state.get(hub_key) or {}
                if state.get("aborted"):
                    raise asyncio.CancelledError()
                await memory_manager.save_message(member["id"], "assistant", reply)
                result = {
                    "conversation_id": str(member["id"]),
                    "role": role,
                    "task": task_text,
                    "reply": reply,
                    "step": step_index,
                }
                results.append(result)
                _append_hub_progress(hub_key, progress_lines, f"- {role} 已返回结果，继续汇总。")
            upstream_outputs = [item for item in results if int(item.get("step", 0)) == step_index]

        hub_reply = await _build_multi_agent_hub_reply(request.message, members, results, plan_text=plan_text)
        await memory_manager.save_message(hub_id, "assistant", hub_reply)
        _runtime_append_event(hub_key, {"type": "done", "content": hub_reply})
        _runtime_finish(hub_key, hub_reply)
        return {"success": True, "data": {"content": hub_reply, "results": results, "plan": plan_text}}
    except asyncio.CancelledError:
        content = _runtime_snapshot(hub_key).get("content", "") or "多Agent任务已被用户终止。"
        content = f"{content}\n\n任务已被用户终止。".strip()
        try:
            await memory_manager.save_message(hub_id, "assistant", content)
        except Exception:
            pass
        _runtime_append_event(hub_key, {"type": "error", "message": "多Agent任务已被用户终止"})
        _runtime_finish(hub_key, content)
        return {"success": True, "data": {"content": content, "results": [], "aborted": True}}
    except HTTPException:
        raise
    except Exception as e:
        _runtime_append_event(hub_key, {"type": "error", "message": str(e)})
        _runtime_finish(hub_key)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        _multi_agent_dispatch_state.pop(hub_key, None)
        unmark_conversation_running(hub_key)


@router.post("/conversations/{conversation_id}/share")
async def share_conversation(conversation_id: int):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation(conversation_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="对话不存在")
        share_id = conversation.get("share_id") or uuid4().hex
        await memory_manager.set_conversation_share_id(conversation_id, share_id)
        return {
            "success": True,
            "data": {
                "share_id": share_id,
                "share_path": f"/share/{share_id}"
            },
            "message": "分享链接已生成"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/share/{share_id}")
async def get_shared_conversation(share_id: str):
    """公开只读分享接口：通过 share_id 返回对话和消息。"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        conversation = await memory_manager.get_conversation_by_share_id(share_id)
        if not conversation:
            raise HTTPException(status_code=404, detail="分享不存在或已失效")

        messages = await memory_manager.get_messages(conversation["id"], limit=1000)
        return {
            "success": True,
            "data": {
                "conversation": conversation,
                "messages": messages,
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/messages")
async def send_message(
    conversation_id: int,
    request: MessageRequest
):
    """发送消息"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        
        # 保存用户消息
        await memory_manager.save_message(
            conversation_id=conversation_id,
            role="user",
            content=request.content
        )
        
        return {
            "success": True,
            "message": "消息已发送"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/conversations/{conversation_id}/messages")
async def get_messages(
    conversation_id: int,
    limit: int = 100
):
    """获取消息历史"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        
        messages = await memory_manager.get_messages(
            conversation_id=conversation_id,
            limit=limit
        )
        
        return {
            "success": True,
            "data": {
                "items": messages,
                "total": len(messages)
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/skills")
async def generate_skill_from_conversation(
    conversation_id: int,
    request: SkillGenerateRequest
):
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        messages = await memory_manager.get_messages(
            conversation_id=conversation_id,
            limit=request.message_limit
        )
        if not messages:
            raise HTTPException(status_code=404, detail="对话不存在或无消息")

        generated = _build_skill_from_conversation(messages, request)
        user_material = "\n\n".join(
            str(item.get("content") or "").strip()
            for item in messages
            if item.get("role") == "user" and str(item.get("content") or "").strip()
        )
        assistant_material = "\n\n".join(
            str(item.get("content") or "").strip()
            for item in messages
            if item.get("role") == "assistant" and str(item.get("content") or "").strip()
        )
        body = await generate_skill_body_from_chat(
            user_message=user_material,
            assistant_response=assistant_material,
            title=generated.get("name") or "对话技能",
            description=generated.get("description") or "",
            opencode_client=opencode_ws if opencode_ws and getattr(opencode_ws, "connected", False) else None,
        )
        skill = get_skill_registry().create_auto_skill(
            name=generated.get("name") or "对话技能",
            description=generated.get("description") or "",
            body=body,
            user_message=user_material,
        )


        return {
            "success": True,
            "data": skill,
            "message": "技能已生成到 Codebot 自动生成技能目录"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── 文件上传与内容提取 ──────────────────────────────────────────────────────

def _extract_text_from_file(filename: str, content_bytes: bytes) -> Optional[str]:
    """
    从文件字节内容提取纯文本。
    支持：txt, md, csv, py, js, ts, json, yaml, html, xml
    以及 docx, xlsx, pptx, pdf（如果安装了对应库）。
    """
    ext = os.path.splitext(filename)[1].lower()

    # ── 纯文本类 ───────────────────────────────────────────────────────
    text_exts = {
        ".txt", ".md", ".csv", ".py", ".js", ".ts", ".jsx", ".tsx",
        ".json", ".yaml", ".yml", ".html", ".htm", ".xml", ".css",
        ".scss", ".less", ".sh", ".bash", ".zsh", ".bat", ".ps1",
        ".sql", ".toml", ".ini", ".cfg", ".conf", ".log", ".rst",
        ".tex", ".r", ".rb", ".go", ".java", ".c", ".cpp", ".h",
        ".hpp", ".cs", ".php", ".swift", ".kt", ".rs", ".dart",
        ".vue", ".svelte",
    }
    if ext in text_exts:
        for enc in ("utf-8", "gbk", "latin-1"):
            try:
                return content_bytes.decode(enc)
            except UnicodeDecodeError:
                continue
        return content_bytes.decode("utf-8", errors="replace")

    # ── Word 文档 ───────────────────────────────────────────────────────
    if ext in (".docx",):
        try:
            from docx import Document as DocxDocument
            import io
            doc = DocxDocument(io.BytesIO(content_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return f"[Word 文档 {filename}，需安装 python-docx 以提取内容]"
        except Exception as e:
            return f"[Word 文档 {filename} 解析失败: {e}]"

    # ── Excel 表格 ──────────────────────────────────────────────────────
    if ext in (".xlsx", ".xls", ".xlsm"):
        try:
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        lines.append("\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(lines) if lines else f"[Excel 文件 {filename} 内容为空]"
        except ImportError:
            return f"[Excel 文件 {filename}，需安装 openpyxl 以提取内容]"
        except Exception as e:
            return f"[Excel 文件 {filename} 解析失败: {e}]"

    # ── PowerPoint ─────────────────────────────────────────────────────
    if ext in (".pptx",):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content_bytes))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## 幻灯片 {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines) if lines else f"[PPT 文件 {filename} 内容为空]"
        except ImportError:
            return f"[PPT 文件 {filename}，需安装 python-pptx 以提取内容]"
        except Exception as e:
            return f"[PPT 文件 {filename} 解析失败: {e}]"

    # ── PDF ─────────────────────────────────────────────────────────────
    if ext in (".pdf",):
        try:
            import pdfplumber
            import io
            with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
                texts = [page.extract_text() or "" for page in pdf.pages]
            text = "\n".join(t for t in texts if t.strip())
            return text if text.strip() else f"[PDF 文件 {filename} 无可提取文本]"
        except ImportError:
            pass
        try:
            import PyPDF2
            import io
            reader = PyPDF2.PdfReader(io.BytesIO(content_bytes))
            texts = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(t for t in texts if t.strip())
            return text if text.strip() else f"[PDF 文件 {filename} 无可提取文本]"
        except ImportError:
            return f"[PDF 文件 {filename}，需安装 pdfplumber 或 PyPDF2 以提取内容]"
        except Exception as e:
            return f"[PDF 文件 {filename} 解析失败: {e}]"

    # ── 图片 ─────────────────────────────────────────────────────────────
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg"):
        # 图片以 base64 编码返回描述信息，不提取文本
        return None  # 由调用方处理图片

    return f"[不支持的文件类型 {ext}，文件名: {filename}]"


def _build_files_context(attached_files: List[AttachedFile]) -> str:
    """
    将附件列表构建为注入 prompt 的上下文字符串。
    """
    if not attached_files:
        return ""

    parts = []
    for f in attached_files:
        if f.is_text:
            # 文本内容直接注入
            content_preview = f.content
            # 限制单文件最多 50000 字符
            if len(content_preview) > 50000:
                content_preview = content_preview[:50000] + "\n... [内容已截断]"
            parts.append(f"【附件：{f.name}】\n```\n{content_preview}\n```")
        else:
            # 二进制附件只注入元信息，避免把整段 base64 回显到聊天文本。
            parts.append(_attachment_summary(f))

    return "\n\n".join(parts)


def _split_search_tokens(query: str) -> List[str]:
    return [part.lower() for part in re.split(r"[\s\u3000]+", (query or "").strip()) if part.strip()]


def _matches_tokens(query: str, *values: str) -> bool:
    tokens = _split_search_tokens(query)
    if not tokens:
        return True
    haystack = " ".join(str(value or "") for value in values).lower()
    return all(token in haystack for token in tokens)


def _configured_knowledge_bases() -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    seen: set[str] = set()

    def add_item(name: str, path_value: str, description: str = "", enabled: bool = True, kb_id: str = ""):
        raw = (path_value or "").strip()
        if not raw:
            return
        try:
            path = Path(raw).expanduser()
            key = str(path.resolve()).lower()
        except Exception:
            path = Path(raw)
            key = str(path).lower()
        if key in seen:
            return
        seen.add(key)
        items.append({
            "id": kb_id or key,
            "name": name or path.name or raw,
            "path": str(path),
            "description": description or "",
            "enabled": bool(enabled),
        })

    obsidian_cfg = getattr(app_config, "obsidian", None)
    if obsidian_cfg:
        add_item("Obsidian Vault", getattr(obsidian_cfg, "vault_path", ""), "Default Obsidian vault")
        for kb in getattr(obsidian_cfg, "knowledge_bases", []) or []:
            data = kb.model_dump() if hasattr(kb, "model_dump") else dict(kb)
            add_item(
                data.get("name") or "",
                data.get("path") or "",
                data.get("description") or "",
                data.get("enabled", True),
                data.get("id") or "",
            )
    return items


def _resolve_knowledge_roots(knowledge_paths: Optional[List[str]]) -> List[Path]:
    configured = _configured_knowledge_bases()
    allowed_by_id = {str(item.get("id")): item for item in configured}
    allowed_by_path = {}
    for item in configured:
        try:
            allowed_by_path[str(Path(item["path"]).expanduser().resolve()).lower()] = item
        except Exception:
            pass

    roots: List[Path] = []
    requested = [str(item).strip() for item in (knowledge_paths or []) if str(item).strip()]
    if not requested:
        requested = [str(item.get("id") or item.get("path")) for item in configured if item.get("enabled", True)]

    for raw in requested:
        item = allowed_by_id.get(raw)
        if not item:
            try:
                item = allowed_by_path.get(str(Path(raw).expanduser().resolve()).lower())
            except Exception:
                item = None
        if not item or not item.get("enabled", True):
            continue
        path = Path(item["path"]).expanduser()
        if path.exists() and path.is_dir():
            roots.append(path.resolve())
    return list(dict.fromkeys(roots))


def _search_markdown_notes(query: str, knowledge_paths: Optional[List[str]], limit: int = 8) -> List[Dict[str, Any]]:
    roots = _resolve_knowledge_roots(knowledge_paths)
    results: List[Dict[str, Any]] = []
    seen: set[str] = set()
    skip_dirs = {".git", ".obsidian", "node_modules", "__pycache__", ".trash", ".venv", "venv"}
    tokens = _split_search_tokens(query)
    for root in roots:
        if len(results) >= limit:
            break
        for path in root.glob("**/*.md"):
            if len(results) >= limit:
                break
            if any(part in skip_dirs for part in path.parts):
                continue
            try:
                key = str(path.resolve()).lower()
                if key in seen:
                    continue
                rel = str(path.relative_to(root)).replace("\\", "/")
                text = path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            title_match = _matches_tokens(query, rel, path.stem)
            content_match = _matches_tokens(query, text[:20000])
            if tokens and not (title_match or content_match):
                continue
            seen.add(key)
            snippet = ""
            if tokens:
                lower_text = text.lower()
                positions = [lower_text.find(token) for token in tokens if lower_text.find(token) >= 0]
                start = max(0, min(positions) - 240) if positions else 0
            else:
                start = 0
            snippet = text[start:start + 1200].strip()
            results.append({
                "name": path.stem,
                "path": str(path),
                "relative_path": rel,
                "root": str(root),
                "snippet": snippet,
            })
    return results


def _build_obsidian_context(message: str, knowledge_paths: Optional[List[str]]) -> str:
    roots = _resolve_knowledge_roots(knowledge_paths)
    if not roots:
        return ""
    notes = _search_markdown_notes(message, knowledge_paths, limit=6)
    lines = [
        "Obsidian mode is active. Treat configured Obsidian vaults and knowledge folders as plain Markdown sources.",
        "Use obsidian-cli when available for Obsidian actions such as search, templates, note creation, moves, and wiki-link-safe operations.",
        "Do not create a vector database for these knowledge bases; preserve and inspect Markdown files directly.",
        "Selected knowledge roots:",
    ]
    lines.extend([f"- {root}" for root in roots])
    if notes:
        lines.append("Relevant Markdown notes from direct text search:")
        for note in notes:
            snippet = (note.get("snippet") or "")[:1200]
            lines.append(f"\n## {note.get('relative_path')}\nPath: {note.get('path')}\n```markdown\n{snippet}\n```")
    return "\n".join(lines)


def _build_hermes_context(model: Optional[str] = None) -> str:
    try:
        from api.routes.hermes import write_bridge_config
        bridge_path = write_bridge_config()
    except Exception:
        bridge_path = settings.DATA_DIR / "hermes" / "codebot_bridge.json"
    background_model = app_config.memory.organize_model or app_config.general.chat_default_model or ""
    active_model = model or app_config.general.chat_default_model or ""
    return (
        "Hermes mode is active. Process this turn as Hermes Agent, with Codebot acting only as the message relay UI.\n"
        "Use Codebot shared resources instead of asking the user to configure separate models.\n"
        f"OpenCode server URL: {app_config.opencode.server_url}\n"
        f"Main chat model: {active_model}\n"
        f"Background memory organization model: {background_model}\n"
        f"Shared memory database: {settings.CONVERSATIONS_DB}\n"
        f"Shared scheduler database: {settings.SCHEDULED_TASKS_DB}\n"
        f"Shared Codebot skills directory: {settings.SKILLS_DIR}\n"
        f"Hermes bridge config: {bridge_path}\n"
        "Hermes skills, memory, scheduled automations, and model availability should be treated as shared with Codebot."
    )


@router.post("/upload_file")
async def upload_file(file: UploadFile = File(...)):
    """
    接收前端上传的文件，提取内容后返回供前端附加到消息中。
    返回格式：{name, type, content, is_text}
    """
    try:
        content_bytes = await file.read()
        filename = file.filename or "unknown"
        mime_type = file.content_type or "application/octet-stream"

        # 尝试提取文本内容
        text_content = _extract_text_from_file(filename, content_bytes)

        if text_content is not None:
            # 文本文件
            return {
                "success": True,
                "data": {
                    "name": filename,
                    "type": mime_type,
                    "content": text_content,
                    "is_text": True,
                }
            }
        else:
            # 图片等二进制文件：base64 编码
            b64 = base64.b64encode(content_bytes).decode("utf-8")
            return {
                "success": True,
                "data": {
                    "name": filename,
                    "type": mime_type,
                    "content": b64,
                    "is_text": False,
                }
            }
    except Exception as e:
        logger.error(f"文件上传处理失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/commands")
async def get_slash_commands():
    """
    返回可用的 / 命令列表（用于前端弹出命令面板）。
    """
    from core.tool_dispatcher import _load_all_skills
    skills = _load_all_skills()

    commands = [
        {
            "name": "skill",
            "label": "/skill",
            "description": "调用一个已安装的技能",
            "icon": "MagicStick",
            "type": "category",
        },
        {
            "name": "clear",
            "label": "/clear",
            "description": "清除当前对话（保留历史）",
            "icon": "Delete",
            "type": "action",
        },
        {
            "name": "memory",
            "label": "/memory",
            "description": "查看与当前话题相关的记忆",
            "icon": "Collection",
            "type": "action",
        },
        {
            "name": "plan",
            "label": "/plan",
            "description": "切换到规划模式（Plan）",
            "icon": "List",
            "type": "action",
        },
        {
            "name": "build",
            "label": "/build",
            "description": "切换到构建模式（Build）",
            "icon": "Tools",
            "type": "action",
        },
        {
            "name": "agent",
            "label": "/agent",
            "description": "切换到智能体模式（Agent）：自我反思与专家协作",
            "icon": "UserFilled",
            "type": "action",
        },
    ]

    # 把技能列表附加为子命令
    skill_commands = []
    for sk in skills:
        skill_commands.append({
            "name": f"skill:{sk['id']}",
            "label": f"/skill {sk['name']}",
            "description": f"[{sk.get('sourceLabel') or sk.get('source_label') or sk.get('source')}] {sk.get('description', '')}",
            "icon": "MagicStick",
            "type": "skill",
            "skill_id": sk["id"],
            "skill_name": sk["name"],
            "skill_source": sk.get("source"),
            "skill_source_label": sk.get("sourceLabel") or sk.get("source_label") or sk.get("source"),
        })

    return {
        "success": True,
        "data": {
            "commands": commands,
            "skills": skill_commands,
        }
    }


def _serialize_skill_for_search(skill: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "id": skill.get("id"),
        "slug": skill.get("slug"),
        "name": skill.get("name"),
        "description": skill.get("description"),
        "source": skill.get("source"),
        "sourceLabel": skill.get("sourceLabel") or skill.get("source_label") or skill.get("source"),
        "path": skill.get("path"),
    }


def _balanced_skill_results(skills: List[Dict[str, Any]], limit: int) -> List[Dict[str, Any]]:
    source_order = [AUTO_GENERATED, BUILTIN, HERMES, OPENCODE, OPENCLAW, EXTERNAL]
    buckets: Dict[str, List[Dict[str, Any]]] = {source: [] for source in source_order}
    buckets.setdefault("other", [])
    for skill in skills:
        source = str(skill.get("source") or "other")
        buckets.setdefault(source, []).append(skill)

    results: List[Dict[str, Any]] = []
    while len(results) < limit:
        added = False
        for source in source_order + ["other"]:
            bucket = buckets.get(source) or []
            if not bucket:
                continue
            results.append(bucket.pop(0))
            added = True
            if len(results) >= limit:
                break
        if not added:
            break
    return results


@router.get("/skills/search")
async def search_chat_skills(query: str = "", limit: int = 50):
    registry = get_skill_registry()
    capped_limit = max(1, min(limit, 100))
    skills = registry.list_skills()
    if not (query or "").strip():
        return {
            "success": True,
            "data": {
                "skills": [_serialize_skill_for_search(skill) for skill in _balanced_skill_results(skills, capped_limit)]
            },
        }

    items = []
    for skill in skills:
        if not _matches_tokens(
            query,
            skill.get("id", ""),
            skill.get("slug", ""),
            skill.get("name", ""),
            skill.get("description", ""),
            skill.get("sourceLabel", ""),
            skill.get("source_label", ""),
        ):
            continue
        items.append(_serialize_skill_for_search(skill))
        if len(items) >= capped_limit:
            break
    return {"success": True, "data": {"skills": items}}


@router.get("/knowledge/search")
async def search_knowledge_bases(query: str = "", limit: int = 20):
    bases = []
    for item in _configured_knowledge_bases():
        if not item.get("enabled", True):
            continue
        if not _matches_tokens(query, item.get("name", ""), item.get("path", ""), item.get("description", "")):
            continue
        bases.append(item)
        if len(bases) >= max(1, min(limit, 50)):
            break
    return {"success": True, "data": {"items": bases}}


@router.get("/files/search")
async def search_files(query: str = "", limit: int = 20, project_dir: str = ""):
    """
    @文件搜索：在工作目录、配置的额外目录以及当前项目目录中搜索文件（用于前端 @ 触发）。
    """
    import glob as glob_module
    from pathlib import Path as _Path

    # 收集所有要搜索的根目录
    search_roots: list[tuple[_Path, str]] = []   # (abs_path, label_prefix)

    # 1) 默认 BASE_DIR
    search_roots.append((settings.BASE_DIR, ""))

    # 2) 配置的额外 file_search_dirs
    try:
        cfg = app_config
        extra_dirs = cfg.general.file_search_dirs or []
        for d in extra_dirs:
            p = _Path(d)
            if p.is_absolute() and p.is_dir():
                search_roots.append((p, f"[{p.name}] "))
    except Exception:
        pass

    # 3) 前端传来的 project_dir
    if project_dir:
        pd = _Path(project_dir)
        if pd.is_absolute() and pd.is_dir():
            # 去重：如果已经在 search_roots 中则跳过
            existing = {str(r[0].resolve()) for r in search_roots}
            if str(pd.resolve()) not in existing:
                search_roots.append((pd, f"[{pd.name}] "))

    results = []
    seen = set()

    try:
        patterns = ["**/*.md", "**/*.txt", "**/*.py", "**/*.js", "**/*.ts",
                    "**/*.json", "**/*.yaml", "**/*.yml", "**/*.csv",
                    "**/*.xlsx", "**/*.docx", "**/*.pdf"]
        skip_dirs = {"node_modules", ".git", "__pycache__", "dist", "build", ".venv", "venv"}

        for root_path, label in search_roots:
            if len(results) >= limit:
                break
            for pattern in patterns:
                if len(results) >= limit:
                    break
                for path in root_path.glob(pattern):
                    parts = path.parts
                    if any(p in skip_dirs for p in parts):
                        continue
                    try:
                        rel = str(path.relative_to(root_path)).replace("\\", "/")
                    except ValueError:
                        continue
                    # 用 absolute path 去重
                    abs_key = str(path.resolve())
                    if abs_key in seen:
                        continue
                    # 关键词过滤
                    if query and query.lower() not in rel.lower():
                        continue
                    seen.add(abs_key)
                    results.append({
                        "path": label + rel,
                        "abs_path": str(path),
                        "name": path.name,
                        "ext": path.suffix.lower(),
                        "root": str(root_path),
                    })
                    if len(results) >= limit:
                        break
    except Exception as e:
        logger.warning(f"文件搜索失败: {e}")

    return {
        "success": True,
        "data": {"files": results[:limit]}
    }


@router.post("/read_file")
async def read_file_content(path: str = Body(..., embed=True), abs_path: str = Body(None, embed=True)):
    """
    读取指定路径的文件内容，供 @ 文件插入使用。
    支持两种方式：
    1. abs_path: 由文件搜索返回的绝对路径（需要在允许的目录范围内）
    2. path: 相对于 BASE_DIR 的相对路径（兼容旧逻辑）
    """
    try:
        from pathlib import Path as _Path

        # 构建允许的目录白名单
        allowed_roots = [settings.BASE_DIR.resolve()]
        try:
            cfg = app_config
            for d in (cfg.general.file_search_dirs or []):
                p = _Path(d)
                if p.is_absolute() and p.is_dir():
                    allowed_roots.append(p.resolve())
        except Exception:
            pass

        # 优先使用 abs_path
        if abs_path:
            full_path = _Path(abs_path).resolve()
        else:
            full_path = (settings.BASE_DIR / path).resolve()

        # 安全检查：必须在允许的目录之一内
        if not any(str(full_path).startswith(str(root)) for root in allowed_roots):
            raise HTTPException(status_code=403, detail="不允许读取此路径")
        if not full_path.exists() or not full_path.is_file():
            raise HTTPException(status_code=404, detail="文件不存在")

        content_bytes = full_path.read_bytes()
        text_content = _extract_text_from_file(full_path.name, content_bytes)
        if text_content is None:
            b64 = base64.b64encode(content_bytes).decode("utf-8")
            return {
                "success": True,
                "data": {
                    "name": full_path.name,
                    "path": path,
                    "content": b64,
                    "is_text": False,
                }
            }
        return {
            "success": True,
            "data": {
                "name": full_path.name,
                "path": path,
                "content": text_content,
                "is_text": True,
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/send")
async def send_to_opencode(request: SendMessageRequest):
    """发送消息到 OpenCode。支持多任务排队：如果该对话已有任务在运行，新任务会加入队列。"""
    conv_id = str(request.conversation_id)

    # 构建包含附件内容的完整消息
    full_message = request.message
    if request.attached_files:
        files_context = _build_files_context(request.attached_files)
        if files_context:
            full_message = f"{files_context}\n\n【用户消息】{request.message}" if request.message.strip() else files_context

    # 如果该对话已有队列，把任务入队并立刻返回"已排队"
    if conv_id not in _task_queues:
        _task_queues[conv_id] = asyncio.Queue()

    if is_conversation_running(conv_id) or not _task_queues[conv_id].empty():
        await _task_queues[conv_id].put({
            "message": full_message,
            "model": request.model,
            "mode": request.mode,
            "project_dir": request.project_dir,
            "target": request.target,
            "knowledge_paths": request.knowledge_paths,
            "user_already_saved": request.user_already_saved,
        })
        return {
            "success": True,
            "data": {"content": None, "queued": True},
            "message": "任务已排队，将在当前任务完成后执行"
        }

    try:
        conversations_db.connect()
        memory_manager = MemoryManager()
        content = await _execute_opencode(
            full_message,
            model=request.model,
            mode=request.mode,
            conversation_id=conv_id,
            project_dir=request.project_dir,
            target=request.target,
            knowledge_paths=request.knowledge_paths,
        )

        if content:
            await memory_manager.save_message(
                conversation_id=request.conversation_id,
                role="assistant",
                content=content
            )
            conversation = await memory_manager.get_conversation(request.conversation_id)
            if conversation:
                existing_title = conversation.get("title") or ""
                if existing_title == "新对话" or existing_title.strip() == "":
                    async def _update_title_bg(conv_id, msg, resp, mdl):
                        try:
                            new_title = await generate_conversation_title_via_ai(msg, resp, model=mdl)
                            await memory_manager.update_conversation_title(conv_id, new_title)
                        except Exception as e:
                            logger.debug(f"后台更新对话标题失败: {e}")
                    asyncio.create_task(_update_title_bg(request.conversation_id, request.message, content, request.model))

            # 后台运行完整学习闭环：记忆、定时任务、技能和成长候选。
            asyncio.create_task(
                _run_chat_post_processing(
                    user_message=request.message,
                    assistant_response=content,
                    conversation_id=request.conversation_id,
                    target=request.target,
                    execution_model=request.model,
                )
            )

        # 处理队列中等待的任务（非阻塞，后台运行）
        asyncio.create_task(_drain_queue(conv_id, request.conversation_id))

        return {
            "success": True,
            "data": {"content": content, "queued": False},
            "message": "消息已处理"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _chunk_text(text: str, chunk_size: int = 24) -> List[str]:
    if not text:
        return []
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


def _redact_attachment_payload(value: Any) -> Any:
    if isinstance(value, list):
        return [_redact_attachment_payload(item) for item in value]
    if not isinstance(value, dict):
        return value

    redacted = {key: _redact_attachment_payload(val) for key, val in value.items()}
    mime_type = str(redacted.get("type") or "")
    content = redacted.get("content")
    if mime_type.startswith("image/") and isinstance(content, str) and content:
        redacted["content"] = f"[omitted image base64, {len(content)} chars]"
    return redacted


def _json_for_display(value: Any) -> str:
    try:
        return json.dumps(_redact_attachment_payload(value), ensure_ascii=False, indent=2, default=str)
    except Exception:
        return str(value)


def _tool_state_text(state: dict, key: str) -> str:
    value = state.get(key) if isinstance(state, dict) else None
    return value if isinstance(value, str) and value.strip() else ""


def _format_tool_part_detail(part: dict) -> str:
    state = part.get("state") if isinstance(part.get("state"), dict) else {}
    lines = [
        f"工具：{part.get('tool') or 'unknown'}",
        f"状态：{state.get('status') or 'unknown'}",
    ]
    title = _tool_state_text(state, "title")
    if title:
        lines.append(f"标题：{title}")
    raw = _tool_state_text(state, "raw")
    if raw:
        lines.extend(["", "原始调用：", raw])
    input_value = state.get("input")
    if input_value:
        lines.extend(["", "输入：", f"```json\n{_json_for_display(input_value)}\n```"])
    output = _tool_state_text(state, "output")
    if output:
        lines.extend(["", "输出：", output])
    error = _tool_state_text(state, "error")
    if error:
        lines.extend(["", "错误：", error])
    metadata = state.get("metadata")
    if metadata:
        lines.extend(["", "元数据：", f"```json\n{_json_for_display(metadata)}\n```"])
    attachments = state.get("attachments")
    if attachments:
        lines.extend(["", "附件：", f"```json\n{_json_for_display(attachments)}\n```"])
    return "\n".join(lines).strip()


def _format_part_detail(part: dict) -> str:
    part_type = part.get("type")
    if part_type == "tool":
        return _format_tool_part_detail(part)
    if part_type == "reasoning":
        return str(part.get("text") or "").strip()
    if part_type == "step-start":
        snapshot = str(part.get("snapshot") or "").strip()
        return f"快照：{snapshot}" if snapshot else ""
    if part_type == "step-finish":
        lines = [f"原因：{part.get('reason') or 'completed'}"]
        tokens = part.get("tokens")
        if tokens:
            lines.extend(["", "Token：", f"```json\n{_json_for_display(tokens)}\n```"])
        if part.get("cost") is not None:
            lines.append(f"费用：{part.get('cost')}")
        return "\n".join(lines)
    if part_type == "patch":
        return f"```json\n{_json_for_display(part.get('files') or part)}\n```"
    if part_type == "snapshot":
        return str(part.get("snapshot") or "").strip()
    if part_type == "agent":
        return f"Agent：{part.get('name') or ''}".strip()
    if part_type == "retry":
        return f"第 {part.get('attempt')} 次重试\n\n```json\n{_json_for_display(part.get('error') or part)}\n```"
    if part_type == "compaction":
        return "上下文已自动压缩" if part.get("auto") else "上下文已压缩"
    return f"```json\n{_json_for_display(part)}\n```"


def _part_to_stream_event(part: dict) -> Optional[dict]:
    if not isinstance(part, dict):
        return None
    part_type = part.get("type")
    if part_type == "text":
        return None
    event_type = part_type or "unknown"
    event_kind = "tool_event" if part_type in {
        "tool",
        "step-start",
        "step-finish",
        "reasoning",
        "subtask",
        "patch",
        "snapshot",
        "agent",
        "retry",
        "compaction",
        "file",
    } else "meta_event"
    summary = ""
    if part_type == "tool":
        state = part.get("state") if isinstance(part.get("state"), dict) else {}
        summary = (
            _tool_state_text(state, "title")
            or f"{part.get('tool') or 'tool'} · {state.get('status') or 'unknown'}"
        )
    elif part_type == "reasoning":
        reasoning_text = str(part.get("text") or "").strip()
        first_line = next((line.strip() for line in reasoning_text.splitlines() if line.strip()), "")
        summary = first_line[:120] or "推理"
    elif part_type == "step-finish":
        summary = str(part.get("reason") or "步骤完成")
    elif part_type == "agent":
        summary = str(part.get("name") or "Agent")
    elif part_type == "subtask":
        summary = str(part.get("description") or part.get("agent") or "子任务")
    else:
        summary = event_type
    return {
        "type": event_kind,
        "event_type": event_type,
        "data": part,
        "summary": summary,
        "detail": _format_part_detail(part),
    }


def _permission_event_to_stream_event(event_type: str, properties: dict) -> dict:
    request_id = properties.get("id") or properties.get("requestID") or properties.get("permissionID")
    permission_name = properties.get("permission") or properties.get("type") or "permission"
    metadata = properties.get("metadata") if isinstance(properties.get("metadata"), dict) else {}
    title = str(properties.get("title") or metadata.get("title") or "").strip()
    patterns = properties.get("patterns") or properties.get("pattern") or []
    if isinstance(patterns, str):
        patterns = [patterns]
    summary = title or f"OpenCode 请求确认：{permission_name}"
    lines = ["OpenCode 正在等待你的选择，任务会在回复前暂停。", f"请求：{permission_name}"]
    if title:
        lines.append(f"标题：{title}")
    if patterns:
        lines.extend(["匹配规则：", *[f"- {item}" for item in patterns]])
    if metadata:
        lines.extend(["", "详细信息：", f"```json\n{_json_for_display(metadata)}\n```"])
    tool = properties.get("tool")
    if isinstance(tool, dict) and tool:
        lines.extend(["", "关联工具调用：", f"```json\n{_json_for_display(tool)}\n```"])
    actions = []
    requires_action = event_type in {"permission.asked", "permission.updated"} and bool(request_id)
    if requires_action:
        options = metadata.get("options") or properties.get("options") or []
        if isinstance(options, list) and options:
            actions = [
                {"label": str(option), "reply": "once", "message": str(option), "type": "primary"}
                for option in options
                if str(option).strip()
            ]
            actions.append({"label": "拒绝", "reply": "reject", "type": "danger"})
        else:
            actions = [
                {"label": "允许一次", "reply": "once", "type": "primary"},
                {"label": "总是允许", "reply": "always", "type": "success"},
                {"label": "拒绝", "reply": "reject", "type": "danger"},
            ]
    if event_type == "permission.replied":
        reply = properties.get("reply") or properties.get("response")
        summary = f"权限请求已回复：{reply or '已处理'}"
    return {
        "type": "meta_event",
        "event_type": event_type,
        "data": properties,
        "summary": summary,
        "detail": "\n".join(lines).strip(),
        "requires_user_action": requires_action,
        "actions": actions,
        "request_id": request_id,
    }


def _question_event_to_stream_event(event_type: str, properties: dict) -> dict:
    request_id = properties.get("id") or properties.get("requestID")
    raw_questions = properties.get("questions") if isinstance(properties.get("questions"), list) else []
    questions = []
    for item in raw_questions:
        if not isinstance(item, dict):
            continue
        raw_options = item.get("options") if isinstance(item.get("options"), list) else []
        options = []
        for option in raw_options:
            if not isinstance(option, dict):
                continue
            label = str(option.get("label") or "").strip()
            if not label:
                continue
            options.append({
                "label": label,
                "description": str(option.get("description") or "").strip(),
            })
        questions.append({
            "header": str(item.get("header") or "Question").strip(),
            "question": str(item.get("question") or "OpenCode 正在等待你的选择").strip(),
            "multiple": bool(item.get("multiple")),
            "custom": item.get("custom") is not False,
            "options": options,
        })
    first = questions[0] if questions else {}
    question_text = str(first.get("question") or "OpenCode 正在等待你的选择").strip()
    header = str(first.get("header") or "Question").strip()
    multiple = bool(first.get("multiple"))
    allow_custom = first.get("custom") is not False
    options = first.get("options") if isinstance(first.get("options"), list) else []
    actions = []
    for idx, option in enumerate(options):
        if not isinstance(option, dict):
            continue
        label = str(option.get("label") or "").strip()
        if not label:
            continue
        actions.append({
            "label": label,
            "reply": "question_answer",
            "answers": [[label]],
            "type": "primary" if idx == 0 else "default",
        })
    if allow_custom:
        actions.append({"label": "自定义回答", "reply": "question_custom", "type": "info", "custom": True})
    actions.append({"label": "取消/先不回答", "reply": "question_reject", "type": "danger"})

    detail_lines = [question_text]
    if options:
        detail_lines.extend(["", "选项："])
        for idx, option in enumerate(options, start=1):
            if not isinstance(option, dict):
                continue
            label = str(option.get("label") or "").strip()
            desc = str(option.get("description") or "").strip()
            if label:
                detail_lines.append(f"{idx}. {label}")
                if desc:
                    detail_lines.append(f"   {desc}")
    if len(questions) > 1:
        detail_lines.append("")
        detail_lines.append(f"还有 {len(questions) - 1} 个问题会在回复后继续处理。")

    return {
        "type": "meta_event",
        "event_type": event_type,
        "data": properties,
        "summary": f"{header}: {question_text}" if header else question_text,
        "detail": "\n".join(detail_lines).strip(),
        "requires_user_action": event_type == "question.asked" and bool(request_id),
        "actions": actions if event_type == "question.asked" and bool(request_id) else [],
        "request_id": request_id,
        "question": question_text,
        "questions": questions,
        "multiple": multiple,
        "allow_custom": allow_custom,
    }


def _opencode_event_to_stream_event(stream_event: dict) -> Optional[dict]:
    event_type = str(stream_event.get("event_type") or "unknown")
    properties = stream_event.get("properties") if isinstance(stream_event.get("properties"), dict) else {}
    if event_type in {"permission.asked", "permission.updated", "permission.replied"}:
        return _permission_event_to_stream_event(event_type, properties)
    if event_type in {"question.asked", "question.replied", "question.rejected"}:
        return _question_event_to_stream_event(event_type, properties)

    summary = ""
    detail = ""
    if event_type == "todo.updated":
        todos = properties.get("todos") if isinstance(properties.get("todos"), list) else []
        summary = f"待办更新：{len(todos)} 项"
        detail = "\n".join(
            f"- [{item.get('status') or 'pending'}] {item.get('content') or ''}"
            for item in todos if isinstance(item, dict)
        )
    elif event_type == "session.status":
        status = properties.get("status")
        status_type = status.get("type") if isinstance(status, dict) else status
        summary = f"会话状态：{status_type or 'updated'}"
        detail = f"```json\n{_json_for_display(status or properties)}\n```"
    elif event_type == "session.idle":
        summary = "会话状态：idle"
    elif event_type == "session.error":
        summary = "OpenCode 会话错误"
        detail = f"```json\n{_json_for_display(properties.get('error') or properties)}\n```"
    elif event_type == "message.updated":
        info = properties.get("info") if isinstance(properties.get("info"), dict) else {}
        if info.get("role") != "assistant":
            return None
        finish = info.get("finish")
        error = info.get("error")
        summary = "助手消息更新"
        if finish:
            summary = f"助手消息完成：{finish}"
        if error:
            summary = "助手消息错误"
        visible = {
            key: info.get(key)
            for key in ["id", "role", "mode", "finish", "cost", "tokens", "error"]
            if info.get(key) is not None
        }
        detail = f"```json\n{_json_for_display(visible)}\n```" if visible else ""
    elif event_type == "message.part.updated":
        part = properties.get("part") if isinstance(properties.get("part"), dict) else {}
        part_type = part.get("type") or "part"
        summary = f"消息片段更新：{part_type}"
        if part_type == "text":
            text = str(part.get("text") or "").strip()
            detail = text or f"```json\n{_json_for_display(part)}\n```"
        else:
            detail = f"```json\n{_json_for_display(part or properties)}\n```"
    elif event_type == "file.edited":
        summary = f"文件已编辑：{properties.get('file') or ''}".strip()
    elif event_type == "command.executed":
        summary = f"命令已执行：{properties.get('name') or ''}".strip()
        detail = f"```json\n{_json_for_display(properties)}\n```"
    elif event_type.startswith("pty."):
        info = properties.get("info") if isinstance(properties.get("info"), dict) else properties
        summary = f"终端事件：{event_type}"
        detail = f"```json\n{_json_for_display(info)}\n```"
    elif event_type.startswith("tui."):
        summary = str(properties.get("message") or properties.get("title") or event_type)
        detail = f"```json\n{_json_for_display(properties)}\n```"
    else:
        summary = event_type
        detail = f"```json\n{_json_for_display(properties)}\n```" if properties else ""

    return {
        "type": "meta_event",
        "event_type": event_type,
        "data": properties,
        "summary": summary,
        "detail": detail,
    }


async def _notify_opencode_action_required(conversation_id: int, event: dict):
    if not event.get("requires_user_action"):
        return
    request_id = str(event.get("request_id") or "")
    key = request_id or hashlib.sha1(_json_for_display(event).encode("utf-8", errors="ignore")).hexdigest()
    if key in _opencode_action_notification_keys:
        return
    _opencode_action_notification_keys.add(key)
    if len(_opencode_action_notification_keys) > 500:
        _opencode_action_notification_keys.clear()
        _opencode_action_notification_keys.add(key)

    title = "OpenCode 等待你的选择"
    summary = str(event.get("summary") or "有一个 OpenCode 操作需要确认")
    message = f"{summary}\n\n对话ID: {conversation_id}\n请回到 Codebot 聊天窗口处理。"
    try:
        from api.routes import notifications as notifications_router
        service = getattr(notifications_router, "notification_service", None)
        if service is None:
            return
        asyncio.create_task(
            service.send_action_required_notification(
                title=title,
                message=message,
                task_id=f"chat:{conversation_id}",
                notif_type="warning",
                force_desktop=True,
            )
        )
    except Exception as exc:
        logger.debug(f"OpenCode 操作提醒发送失败（跳过）: {exc}")


def _opencode_cli_display_enabled() -> bool:
    return bool(getattr(app_config.general, "opencode_cli_display", True))


def _cli_tool_symbol(tool_name: str, title: str = "") -> str:
    name = (tool_name or title or "").lower()
    if "webfetch" in name or "web_fetch" in name:
        return "%"
    if "glob" in name:
        return "*"
    if "bash" in name or "shell" in name:
        return "$"
    if any(token in name for token in ["edit", "write", "patch"]):
        return "✎"
    if "task" in name:
        return "◌"
    return "→"


def _cli_todo_marker(status: str) -> str:
    normalized = (status or "").lower()
    if normalized in {"completed", "done", "success"}:
        return "[✓]"
    if normalized in {"in_progress", "running", "active"}:
        return "[•]"
    if normalized in {"cancelled", "canceled", "error", "failed"}:
        return "[-]"
    return "[ ]"


def _cli_append_block(current: str, block: str) -> Tuple[str, str]:
    text = (block or "").strip("\n")
    if not text:
        return current, ""
    prefix = "\n\n" if current and not current.endswith("\n\n") else ""
    delta = f"{prefix}{text}\n"
    return f"{current}{delta}", delta


def _tool_dict(value: Any) -> Dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _tool_path(value: Any) -> str:
    if not isinstance(value, str) or not value.strip():
        return ""
    raw = value.strip()
    try:
        path_obj = Path(raw)
        if path_obj.is_absolute():
            try:
                rel = path_obj.relative_to(Path.cwd())
                return str(rel).replace("\\", "/") or "."
            except ValueError:
                return str(path_obj).replace("\\", "/")
    except Exception:
        pass
    return raw.replace("\\", "/")


def _cli_info(data: Dict[str, Any], skip: Optional[set[str]] = None) -> str:
    skip = skip or set()
    items = []
    for key, val in data.items():
        if key in skip:
            continue
        if isinstance(val, (str, int, float, bool)) and str(val).strip():
            items.append(f"{key}={val}")
    return f"[{', '.join(items)}]" if items else ""


def _cli_count(value: Any, label: str) -> str:
    try:
        count_value = int(value)
    except Exception:
        return ""
    return f"{count_value} {label}{'' if count_value == 1 else 'es'}"


def _format_cli_tool_title(tool: str, state: Dict[str, Any], input_value: Dict[str, Any], metadata: Dict[str, Any]) -> Tuple[str, str, bool]:
    name = (tool or "").lower()
    explicit_title = _tool_state_text(state, "title")
    if name in {"read", "view"}:
        file_path = _tool_path(input_value.get("filePath") or input_value.get("file_path") or input_value.get("path"))
        extra = _cli_info(input_value, {"filePath", "file_path", "path"})
        base = f"Read {file_path}" if file_path else (explicit_title if explicit_title.lower().startswith("read") else f"Read {explicit_title}".strip())
        return "→", f"{base}{f' {extra}' if extra else ''}".strip(), False
    if name == "skill":
        skill_name = str(input_value.get("name") or "").strip()
        base = f"Skill \"{skill_name}\"" if skill_name else (explicit_title if explicit_title.lower().startswith("skill") else f"Skill \"{explicit_title}\"")
        return "→", base.strip(), False
    if name == "webfetch":
        url = str(input_value.get("url") or "").strip()
        return "%", f"WebFetch {url}".strip() or "WebFetch", False
    if name == "websearch":
        query = str(input_value.get("query") or "").strip()
        return "◈", f"WebSearch \"{query}\"" if query else "WebSearch", False
    if name == "glob":
        pattern = str(input_value.get("pattern") or "").strip()
        root = _tool_path(input_value.get("path"))
        matches = _cli_count(metadata.get("count"), "match")
        desc = f" in {root}" if root else ""
        desc = f"{desc} · {matches}" if matches and desc else (f" · {matches}" if matches else desc)
        return "✱", f"Glob \"{pattern}\"{desc}".strip(), False
    if name == "grep":
        pattern = str(input_value.get("pattern") or "").strip()
        root = _tool_path(input_value.get("path"))
        matches = _cli_count(metadata.get("matches"), "match")
        desc = f" in {root}" if root else ""
        desc = f"{desc} · {matches}" if matches and desc else (f" · {matches}" if matches else desc)
        return "✱", f"Grep \"{pattern}\"{desc}".strip(), False
    if name == "list":
        root = _tool_path(input_value.get("path"))
        return "→", f"List {root}".strip(), False
    if name == "lsp":
        op = input_value.get("operation") or "request"
        file_path = _tool_path(input_value.get("filePath") or input_value.get("file_path"))
        line = input_value.get("line")
        char = input_value.get("character")
        pos = f":{line}:{char}" if isinstance(line, int) and isinstance(char, int) else ""
        return "→", f"LSP {op} {file_path}{pos}".strip(), False
    if name == "todowrite":
        return "#", "Todos", False
    if name == "question":
        total = len(input_value.get("questions") or []) if isinstance(input_value.get("questions"), list) else 0
        return "→", f"Asked {total} question{'' if total == 1 else 's'}", False
    if name == "task":
        kind = str(input_value.get("subagent_type") or "general").title()
        desc = str(input_value.get("description") or explicit_title or "").strip()
        status = str(state.get("status") or "").lower()
        icon = "✗" if status == "error" else ("•" if status == "running" else "✓")
        return icon, desc or f"{kind} Task", False
    if name in {"edit", "write"}:
        file_path = _tool_path(input_value.get("filePath") or input_value.get("file_path") or input_value.get("path"))
        return "←", f"{name.title()} {file_path}".strip(), True
    if name == "apply_patch":
        files = metadata.get("files")
        count_value = len(files) if isinstance(files, list) else 0
        title = "Patch" if count_value == 0 else f"Patch {count_value} file{'' if count_value == 1 else 's'}"
        return "%", title, False
    if name in {"bash", "shell"}:
        command = str(input_value.get("command") or explicit_title or "").strip()
        return "$", command or "Bash", True
    return _cli_tool_symbol(tool, explicit_title), explicit_title or tool or "Tool", False


def _format_todo_body_from_input(input_value: Dict[str, Any]) -> str:
    todos = input_value.get("todos")
    if not isinstance(todos, list):
        return ""
    rows = []
    for item in todos:
        if not isinstance(item, dict):
            continue
        content = str(item.get("content") or "").strip()
        if content:
            rows.append(f"{_cli_todo_marker(str(item.get('status') or ''))} {content}")
    return "\n".join(rows)


def _format_cli_tool_part(part: dict, seen: set[str]) -> str:
    state = _tool_dict(part.get("state"))
    input_value = _tool_dict(state.get("input"))
    metadata = _tool_dict(state.get("metadata"))
    part_id = str(part.get("id") or part.get("callID") or "")
    tool = str(part.get("tool") or "")
    icon, title, allow_body = _format_cli_tool_title(tool, state, input_value, metadata)
    status = str(state.get("status") or "").lower()
    error = _tool_state_text(state, "error")
    if status == "error" and error:
        title = f"{title} failed: {error}"
        icon = "✗"

    lines: List[str] = []
    title_key = f"tool-title:{part_id or tool}:{icon}:{title}"
    if title and title_key not in seen:
        seen.add(title_key)
        lines.append(f"{icon} {title}".strip())

    body = ""
    name = (tool or "").lower()
    if name == "todowrite":
        body = _format_todo_body_from_input(input_value)
    elif allow_body:
        if name == "edit":
            body = str(metadata.get("diff") or "").strip()
        else:
            body = _tool_state_text(state, "output").strip()
    if body:
        body_key = f"tool-body:{part_id}:{hashlib.sha1(body.encode('utf-8', errors='ignore')).hexdigest()}"
        if body_key not in seen:
            seen.add(body_key)
            if lines:
                lines.append("")
            lines.append(body.rstrip())

    return "\n".join(lines).strip()


def _format_cli_part(part: dict, seen: set[str]) -> str:
    if not isinstance(part, dict):
        return ""
    part_type = part.get("type")
    if part_type == "tool":
        return _format_cli_tool_part(part, seen)
    if part_type == "reasoning":
        return ""
    if part_type == "agent":
        return ""
    if part_type == "patch":
        files = part.get("files") if isinstance(part.get("files"), list) else []
        if not files:
            return ""
        key = f"patch:{part.get('id')}:{len(files)}"
        if key in seen:
            return ""
        seen.add(key)
        return "\n".join(f"✎ {item}" for item in files)
    return ""


def _format_cli_opencode_event(stream_event: dict, seen: set[str]) -> str:
    event_type = str(stream_event.get("event_type") or "")
    properties = stream_event.get("properties") if isinstance(stream_event.get("properties"), dict) else {}

    if event_type == "todo.updated":
        todos = properties.get("todos") if isinstance(properties.get("todos"), list) else []
        lines = ["# Todos", ""]
        for item in todos:
            if not isinstance(item, dict):
                continue
            lines.append(f"{_cli_todo_marker(str(item.get('status') or ''))} {item.get('content') or ''}".rstrip())
        block = "\n".join(lines).strip()
        key = f"todo:{hashlib.sha1(block.encode('utf-8', errors='ignore')).hexdigest()}"
        if key in seen:
            return ""
        seen.add(key)
        return block

    if event_type in {"permission.asked", "permission.updated"}:
        converted = _permission_event_to_stream_event(event_type, properties)
        request_id = str(converted.get("request_id") or "")
        key = f"permission:{request_id}"
        if key in seen:
            return ""
        seen.add(key)
        detail = str(converted.get("summary") or "OpenCode 正在等待你的选择").strip()
        actions = converted.get("actions") if isinstance(converted.get("actions"), list) else []
        action_text = " / ".join(str(action.get("label") or "") for action in actions if isinstance(action, dict))
        return f"! {detail}\n{action_text}".strip()

    if event_type == "question.asked":
        converted = _question_event_to_stream_event(event_type, properties)
        request_id = str(converted.get("request_id") or "")
        key = f"question:{request_id}"
        if key in seen:
            return ""
        seen.add(key)
        lines = [f"! {converted.get('question') or converted.get('summary') or 'OpenCode 正在等待你的选择'}"]
        question_items = properties.get("questions") if isinstance(properties.get("questions"), list) else []
        first_question = question_items[0] if question_items and isinstance(question_items[0], dict) else {}
        options = first_question.get("options") if isinstance(first_question.get("options"), list) else []
        if isinstance(options, list):
            for idx, option in enumerate(options, start=1):
                if isinstance(option, dict) and option.get("label"):
                    lines.append(f"{idx}. {option.get('label')}")
        lines.append("请在聊天窗口中选择。")
        return "\n".join(lines).strip()

    if event_type in {"question.replied", "question.rejected"}:
        key = f"{event_type}:{properties.get('requestID') or properties.get('id')}"
        if key in seen:
            return ""
        seen.add(key)
        if event_type == "question.rejected":
            return "✓ Question dismissed"
        answers = properties.get("answers")
        return f"✓ Question answered {answers or ''}".strip()

    if event_type == "permission.replied":
        reply = properties.get("reply") or properties.get("response")
        key = f"permission-replied:{properties.get('requestID') or properties.get('permissionID')}:{reply}"
        if key in seen:
            return ""
        seen.add(key)
        return f"✓ Permission {reply or 'replied'}"

    if event_type == "session.error":
        block = f"! OpenCode error\n{_json_for_display(properties.get('error') or properties)}"
        key = f"session-error:{hashlib.sha1(block.encode('utf-8', errors='ignore')).hexdigest()}"
        if key in seen:
            return ""
        seen.add(key)
        return block

    if event_type == "file.edited":
        file_path = str(properties.get("file") or "").strip()
        key = f"file-edited:{file_path}"
        if not file_path or key in seen:
            return ""
        seen.add(key)
        return f"✎ {file_path}"

    if event_type == "command.executed":
        name = str(properties.get("name") or "").strip()
        args = str(properties.get("arguments") or "").strip()
        block = f"> {name} {args}".strip()
        key = f"command:{hashlib.sha1(block.encode('utf-8', errors='ignore')).hexdigest()}"
        if not name or key in seen:
            return ""
        seen.add(key)
        return block

    return ""


@router.post("/send_stream")
async def send_to_opencode_stream(request: SendMessageRequest):
    conv_id = str(request.conversation_id)
    full_message = request.message
    if request.attached_files:
        files_context = _build_files_context(request.attached_files)
        if files_context:
            full_message = f"{files_context}\n\n【用户消息】{request.message}" if request.message.strip() else files_context

    if conv_id not in _task_queues:
        _task_queues[conv_id] = asyncio.Queue()

    if is_conversation_running(conv_id) or not _task_queues[conv_id].empty():
        await _task_queues[conv_id].put({
            "message": full_message,
            "model": request.model,
            "mode": request.mode,
            "project_dir": request.project_dir,
            "target": request.target,
            "knowledge_paths": request.knowledge_paths,
            "user_already_saved": request.user_already_saved,
        })

        async def queued_stream():
            event = {"type": "queued", "message": "任务已排队，将在当前任务完成后执行"}
            yield json.dumps(event, ensure_ascii=False) + "\n"
        return StreamingResponse(
            queued_stream(),
            media_type="application/x-ndjson",
            headers={
                "Cache-Control": "no-cache, no-transform",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive",
            }
        )

    async def _run_stream_worker(event_queue: asyncio.Queue):
        _runtime_start(conv_id)
        try:
            conversations_db.connect()
            memory_manager = MemoryManager()
            await event_queue.put({"type": "status", "phase": "started"})
            _runtime_append_event(conv_id, {"type": "status", "phase": "started"})
            raw_content = ""
            content = ""
            parts: List[dict] = []
            internal_prompt: str = ""
            tool_events_log: List[dict] = []
            cli_display = _opencode_cli_display_enabled()
            cli_seen: set[str] = set()
            cli_text_started = False
            async for stream_event in _stream_execute_opencode_with_meta(
                full_message,
                model=request.model,
                mode=request.mode,
                conversation_id=conv_id,
                project_dir=request.project_dir,
                target=request.target,
                knowledge_paths=request.knowledge_paths,
            ):
                event_type = stream_event.get("type")
                if event_type == "internal_prompt":
                    internal_prompt = stream_event.get("prompt") or ""
                    continue
                if event_type == "content_delta":
                    delta_text = stream_event.get("delta") or ""
                    raw_content = stream_event.get("content") or f"{raw_content}{delta_text}"
                    if cli_display:
                        if not delta_text:
                            continue
                        prefix = ""
                        if not cli_text_started and content and not content.endswith("\n\n"):
                            prefix = "\n\n"
                        cli_text_started = True
                        delta = f"{prefix}{delta_text}" if prefix else delta_text
                        content = f"{content}{delta}"
                        _runtime_set_content(conv_id, content)
                        await event_queue.put({
                            "type": "content_delta",
                            "delta": delta,
                            "content": content,
                            "cli_display": True,
                        })
                        continue
                    next_content = _sanitize_assistant_output(raw_content, user_message=request.message)
                    if next_content == content:
                        continue
                    delta = next_content[len(content):] if next_content.startswith(content) else next_content
                    content = next_content
                    _runtime_set_content(conv_id, content)
                    await event_queue.put({
                        "type": "content_delta",
                        "delta": delta,
                        "content": content
                    })
                    continue
                if event_type == "tool_event":
                    part = stream_event.get("part")
                    if isinstance(part, dict):
                        parts.append(part)
                        tool_events_log.append(part)
                        if cli_display:
                            block = _format_cli_part(part, cli_seen)
                            content, delta = _cli_append_block(content, block)
                            if delta:
                                _runtime_set_content(conv_id, content)
                                await event_queue.put({
                                    "type": "content_delta",
                                    "delta": delta,
                                    "content": content,
                                    "cli_display": True,
                                })
                            continue
                        converted = _part_to_stream_event(part)
                        if converted is not None:
                            _runtime_append_event(conv_id, converted)
                            await event_queue.put(converted)
                            await _notify_opencode_action_required(request.conversation_id, converted)
                    continue
                if event_type == "opencode_event":
                    converted = _opencode_event_to_stream_event(stream_event)
                    if cli_display:
                        block = _format_cli_opencode_event(stream_event, cli_seen)
                        content, delta = _cli_append_block(content, block)
                        if delta:
                            _runtime_set_content(conv_id, content)
                            await event_queue.put({
                                "type": "content_delta",
                                "delta": delta,
                                "content": content,
                                "cli_display": True,
                            })
                        if converted is not None:
                            tool_events_log.append(converted)
                            if converted.get("requires_user_action"):
                                inline_event = {**converted, "cli_inline": True}
                                _runtime_append_event(conv_id, inline_event)
                                await event_queue.put(inline_event)
                                await _notify_opencode_action_required(request.conversation_id, inline_event)
                        continue
                    if converted is not None:
                        tool_events_log.append(converted)
                        _runtime_append_event(conv_id, converted)
                        await event_queue.put(converted)
                        await _notify_opencode_action_required(request.conversation_id, converted)
                    continue
                if event_type == "done":
                    raw_content = stream_event.get("content") or raw_content
                    if cli_display:
                        content = content or raw_content
                    else:
                        content = _sanitize_assistant_output(raw_content, user_message=request.message) or content
                    _runtime_set_content(conv_id, content)
                    stream_parts = stream_event.get("parts")
                    if isinstance(stream_parts, list):
                        for p in stream_parts:
                            if isinstance(p, dict):
                                parts.append(p)
                                if cli_display:
                                    block = _format_cli_part(p, cli_seen)
                                    content, delta = _cli_append_block(content, block)
                                    if delta:
                                        _runtime_set_content(conv_id, content)
                                        await event_queue.put({
                                            "type": "content_delta",
                                            "delta": delta,
                                            "content": content,
                                            "cli_display": True,
                                        })
                    continue
                if event_type == "error":
                    raise RuntimeError(stream_event.get("error") or "OpenCode 流式调用失败")

            # opencode agent 已经原生处理了所有工具调用（包括 MCP、bash、git 等）。
            # codebot 只负责展示结果，不再做任何二次工具调用或总结。

            if content:
                await memory_manager.save_message(
                    conversation_id=request.conversation_id,
                    role="assistant",
                    content=content
                )
                conversation = await memory_manager.get_conversation(request.conversation_id)
                if conversation:
                    existing_title = conversation.get("title") or ""
                    if existing_title == "新对话" or existing_title.strip() == "":
                        async def _update_title_bg_stream(conv_id, msg, resp, mdl):
                            try:
                                new_title = await generate_conversation_title_via_ai(msg, resp, model=mdl)
                                await memory_manager.update_conversation_title(conv_id, new_title)
                            except Exception as e:
                                logger.debug(f"后台更新对话标题失败(stream): {e}")
                        asyncio.create_task(_update_title_bg_stream(request.conversation_id, request.message, content, request.model))

                asyncio.create_task(
                    _run_chat_post_processing(
                        user_message=request.message,
                        assistant_response=content,
                        conversation_id=request.conversation_id,
                        target=request.target,
                        execution_model=request.model,
                    )
                )

            # 保存聊天日志（内部提示词 + 推理过程 + 最终回复）
            try:
                _save_chat_log(
                    conversation_id=request.conversation_id,
                    user_message=request.message,
                    internal_prompt=internal_prompt,
                    tool_events=tool_events_log,
                    final_reply=content,
                    model=request.model,
                    mode=request.mode,
                )
            except Exception as _log_err:
                logger.warning(f"聊天日志保存失败（跳过）: {_log_err}")

            asyncio.create_task(_drain_queue(conv_id, request.conversation_id))
            _runtime_append_event(conv_id, {"type": "done", "content": content or "", "cli_display": cli_display})
            await event_queue.put({"type": "done", "content": content or "", "cli_display": cli_display})
        except Exception as e:
            _runtime_append_event(conv_id, {"type": "error", "message": str(e)})
            await event_queue.put({"type": "error", "message": str(e)})
        finally:
            _runtime_finish(conv_id, content)
            await event_queue.put({"type": "__worker_done__"})

    async def event_stream():
        event_queue: asyncio.Queue = asyncio.Queue()
        worker_task = asyncio.create_task(_run_stream_worker(event_queue))
        try:
            while True:
                event = await event_queue.get()
                if event.get("type") == "__worker_done__":
                    break
                yield json.dumps(event, ensure_ascii=False) + "\n"
        except asyncio.CancelledError:
            logger.info(f"对话 {conv_id} 的流式连接已断开，任务继续在后台执行")
            return
        finally:
            if worker_task.done():
                try:
                    _ = worker_task.result()
                except Exception:
                    pass

    return StreamingResponse(
        event_stream(),
        media_type="application/x-ndjson",
        headers={
            "Cache-Control": "no-cache, no-transform",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        }
    )


@router.post("/permission/reply")
async def reply_opencode_permission(request: PermissionReplyRequest):
    request_id = (request.request_id or "").strip()
    reply = (request.reply or "").strip()
    if not request_id:
        raise HTTPException(status_code=400, detail="缺少权限请求 ID")
    if reply not in {"once", "always", "reject"}:
        raise HTTPException(status_code=400, detail="无效的权限回复")

    client = opencode_ws or OpenCodeClient(app_config.opencode.server_url)
    ok = await _ensure_opencode_client_connected(client)
    if not ok:
        raise HTTPException(status_code=503, detail="OpenCode 未连接")

    success = await client.reply_permission(
        request_id=request_id,
        reply=reply,
        message=request.message,
        session_id=request.session_id,
        workspace=request.project_dir,
    )
    if not success:
        raise HTTPException(status_code=502, detail="OpenCode 权限回复失败")

    if request.conversation_id:
        _runtime_append_event(str(request.conversation_id), {
            "type": "meta_event",
            "event_type": "permission.local_reply",
            "summary": f"你已选择：{reply}",
            "detail": "",
            "data": {"request_id": request_id, "reply": reply},
        })
    return {"success": True, "message": "已回复 OpenCode 权限请求"}


@router.post("/question/reply")
async def reply_opencode_question(request: QuestionReplyRequest):
    request_id = (request.request_id or "").strip()
    if not request_id:
        raise HTTPException(status_code=400, detail="缺少 question 请求 ID")

    client = opencode_ws or OpenCodeClient(app_config.opencode.server_url)
    ok = await _ensure_opencode_client_connected(client)
    if not ok:
        raise HTTPException(status_code=503, detail="OpenCode 未连接")

    if request.reject:
        success = await client.reject_question(request_id=request_id, workspace=request.project_dir)
        if not success:
            raise HTTPException(status_code=502, detail="拒绝 OpenCode 问题失败")
        reply_text = "已取消/先不回答"
    else:
        answers = request.answers
        if answers is None:
            answer = (request.answer or "").strip()
            if not answer:
                raise HTTPException(status_code=400, detail="缺少问题回答")
            answers = [[answer]]
        cleaned_answers = []
        for answer_group in answers:
            if isinstance(answer_group, list):
                cleaned_answers.append([str(item).strip() for item in answer_group if str(item).strip()])
            else:
                cleaned_answers.append([])
        if not any(cleaned_answers):
            raise HTTPException(status_code=400, detail="缺少有效问题回答")
        success = await client.reply_question(
            request_id=request_id,
            answers=cleaned_answers,
            workspace=request.project_dir,
        )
        if not success:
            raise HTTPException(status_code=502, detail="回复 OpenCode 问题失败")
        reply_text = "；".join(", ".join(group) for group in cleaned_answers)

    if request.conversation_id:
        _runtime_append_event(str(request.conversation_id), {
            "type": "meta_event",
            "event_type": "question.local_reply",
            "summary": f"你已回复：{reply_text}",
            "detail": "",
            "data": {"request_id": request_id, "reply": reply_text, "rejected": request.reject},
        })
    return {"success": True, "message": "已回复 OpenCode 问题"}


async def _drain_queue(conv_id: str, conversation_id: int):
    """处理队列中排队的任务"""
    if conv_id not in _task_queues:
        return
    q = _task_queues[conv_id]
    while not q.empty():
        try:
            task = await asyncio.wait_for(q.get(), timeout=0.1)
        except asyncio.TimeoutError:
            break
        try:
            conversations_db.connect()
            memory_manager = MemoryManager()

            # 前端已保存的排队消息不重复入库；直连 API 入队的消息仍在这里保存。
            if not task.get("user_already_saved"):
                await memory_manager.save_message(
                    conversation_id=conversation_id,
                    role="user",
                    content=task["message"]
                )

            content = await _execute_opencode(
                task["message"],
                model=task.get("model"),
                mode=task.get("mode"),
                conversation_id=conv_id,
                project_dir=task.get("project_dir"),
                target=task.get("target"),
                knowledge_paths=task.get("knowledge_paths"),
            )
            if content:
                await memory_manager.save_message(
                    conversation_id=conversation_id,
                    role="assistant",
                    content=content
                )
                asyncio.create_task(
                    _run_chat_post_processing(
                        user_message=task["message"],
                        assistant_response=content,
                        conversation_id=conversation_id,
                        target=task.get("target"),
                        execution_model=task.get("model"),
                    )
                )
        except Exception as e:
            logger.error(f"处理排队任务失败: {e}")
        finally:
            q.task_done()


class AbortRequest(BaseModel):
    conversation_id: int


@router.post("/abort")
async def abort_task(request: AbortRequest):
    """终止指定对话的当前运行任务"""
    conv_id = str(request.conversation_id)
    client = opencode_ws
    target_conv_ids = [conv_id]
    dispatch_state = _multi_agent_dispatch_state.get(conv_id)
    if dispatch_state:
        dispatch_state["aborted"] = True
        dispatch_state["updated_at"] = datetime.now()
        for member_id in dispatch_state.get("member_ids") or []:
            member_key = str(member_id)
            if member_key and member_key not in target_conv_ids:
                target_conv_ids.append(member_key)

    # 清空队列中等待的任务
    for target_id in target_conv_ids:
        if target_id not in _task_queues:
            continue
        q = _task_queues[target_id]
        cleared = 0
        while not q.empty():
            try:
                q.get_nowait()
                q.task_done()
                cleared += 1
            except Exception:
                break
        if cleared:
            logger.info(f"已清空对话 {target_id} 的 {cleared} 个排队任务")

    # 终止当前正在运行的 OpenCode session
    aborted_sessions = 0
    aborted_hermes = 0
    for target_id in target_conv_ids:
        try:
            from api.routes import hermes as hermes_router
            if await hermes_router.abort_hermes_conversation(target_id):
                aborted_hermes += 1
                logger.info(f"已终止对话 {target_id} 的 Hermes CLI 进程")
        except Exception as e:
            logger.warning(f"终止 Hermes CLI 出错: {e}")

        session_id = _conversation_current_session.get(target_id)
        if not session_id or not client:
            unmark_conversation_running(target_id)
            continue
        try:
            ok = await client.abort_session(session_id)
            if ok:
                aborted_sessions += 1
            logger.info(f"终止对话 {target_id} session {session_id}: {'成功' if ok else '失败'}")
        except Exception as e:
            logger.warning(f"终止 session 出错: {e}")
        finally:
            _conversation_current_session.pop(target_id, None)
            unmark_conversation_running(target_id)

    _runtime_append_event(conv_id, {"type": "error", "message": "任务已被用户终止"})
    _runtime_finish(conv_id)

    return {
        "success": True,
        "message": "已发送终止信号",
        "data": {"conversations": target_conv_ids, "aborted_sessions": aborted_sessions, "aborted_hermes": aborted_hermes}
    }


@router.get("/queue_status/{conversation_id}")
async def get_queue_status(conversation_id: int, since_seq: int = 0):
    """获取对话的任务队列状态"""
    conv_id = str(conversation_id)
    queue_size = 0
    if conv_id in _task_queues:
        queue_size = _task_queues[conv_id].qsize()
    has_running = is_conversation_running(conv_id)
    runtime = _runtime_snapshot(conv_id, since_seq=since_seq)
    return {
        "success": True,
        "data": {
            "running": has_running,
            "queued": queue_size,
            "runtime_events": runtime.get("events", []),
            "runtime_last_seq": runtime.get("last_seq", int(since_seq)),
            "runtime_content": runtime.get("content", ""),
        }
    }


class UndoMessageRequest(BaseModel):
    message_id: int
    conversation_id: int


@router.post("/conversations/{conversation_id}/undo")
async def undo_message(conversation_id: int, request: UndoMessageRequest):
    """撤销指定消息及其之后的所有消息"""
    try:
        conversations_db.connect()
        memory_manager = MemoryManager()

        # 获取所有消息
        messages = await memory_manager.get_messages(conversation_id=conversation_id, limit=1000)
        # 找到目标消息的位置
        target_idx = None
        for i, msg in enumerate(messages):
            if msg.get("id") == request.message_id:
                target_idx = i
                break

        if target_idx is None:
            raise HTTPException(status_code=404, detail="消息不存在")

        # 删除该消息及其之后的所有消息
        to_delete = messages[target_idx:]
        deleted_count = 0
        for msg in to_delete:
            msg_id = msg.get("id")
            if msg_id:
                try:
                    await memory_manager.delete_message(msg_id)
                    deleted_count += 1
                except Exception as e:
                    logger.warning(f"删除消息 {msg_id} 失败: {e}")

        return {
            "success": True,
            "data": {"deleted_count": deleted_count},
            "message": f"已撤销 {deleted_count} 条消息"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/models")
async def get_models():
    """获取 OpenCode 可用模型列表"""
    global opencode_ws
    client = opencode_ws
    created_client = False
    try:
        if client is None:
            client = OpenCodeClient(app_config.opencode.server_url)
            created_client = True
        ok = await _ensure_opencode_client_connected(client)
        if not ok:
            return {"success": False, "data": {"models": []}, "message": "OpenCode 未连接"}
        models = await client.get_models()
        return {"success": True, "data": {"models": models}, "message": ""}
    except Exception as e:
        logger.error(f"获取模型列表失败: {e}")
        return {"success": False, "data": {"models": []}, "message": str(e)}
    finally:
        if created_client and client and getattr(client, "connected", False):
            try:
                await client.disconnect()
            except Exception:
                pass


@router.get("/agents")
async def get_agents():
    """获取 OpenCode 支持的 agent 模式列表"""
    return {
        "success": True,
        "data": {
            "agents": [
                {"id": "plan", "name": "Plan 模式", "description": "规划模式：分析需求、制定计划，适合复杂任务拆解"},
                {"id": "build", "name": "Build 模式", "description": "构建模式：直接执行编码任务，适合快速实现"}
            ]
        },
        "message": ""
    }
